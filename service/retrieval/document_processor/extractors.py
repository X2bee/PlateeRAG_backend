"""
파일별 텍스트 추출 모듈
"""

import logging
import aiofiles
from pathlib import Path
from typing import List, Dict, Any
import PyPDF2
from docx import Document

from .dependencies import (
    PDFMINER_AVAILABLE, XLRD_AVAILABLE, OPENPYXL_AVAILABLE, PYTHON_PPTX_AVAILABLE
)
from .text_utils import TextUtils
from .constants import ENCODINGS

if PDFMINER_AVAILABLE:
    from pdfminer.high_level import extract_text

if OPENPYXL_AVAILABLE and XLRD_AVAILABLE :
    from openpyxl import load_workbook
    import xlrd

if PYTHON_PPTX_AVAILABLE:
    from pptx import Presentation

logger = logging.getLogger("document-processor")

class DocumentExtractor:
    """문서 텍스트 추출 클래스"""
    
    def __init__(self, ocr_processor, config_manager):
        self.ocr_processor = ocr_processor
        self.config_manager = config_manager
    
    # PDF 관련 메서드들
    async def extract_text_from_pdf(self, file_path: str) -> str:
        """PDF 파일에서 텍스트 추출 (실시간 설정에 따라 텍스트 추출 또는 OCR)"""
        try:
            current_config = self.config_manager.get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            logger.info(f"🔄 Real-time PDF processing with provider: {provider}")
            
            # no_model인 경우에만 기본 텍스트 추출
            if provider == 'no_model':
                logger.info("Using text extraction mode (no_model)")
                
                # 1단계: pdfminer 시도
                if PDFMINER_AVAILABLE:
                    logger.info(f"Using pdfminer for {file_path}")
                    try:
                        text = extract_text(file_path)
                        cleaned_text = TextUtils.clean_text(text)
                        if len(cleaned_text.strip()) > 100:
                            logger.info(f"Text extracted via pdfminer: {len(cleaned_text)} chars")
                            return cleaned_text
                    except Exception as e:
                        logger.warning(f"pdfminer failed: {e}")
                        
                # 2단계: PyPDF2 fallback
                logger.info(f"Using PyPDF2 fallback for {file_path}")
                text = await self._extract_text_from_pdf_fallback(file_path)
                logger.info(f"Text extracted via PyPDF2: {len(text)} chars")
                return text
            
            else:
                # openai, vllm 등 다른 프로바이더인 경우 무조건 OCR
                logger.info(f"Using OCR mode with provider: {provider}")
                return await self._extract_text_from_pdf_via_ocr(file_path)
                
        except Exception as e:
            logger.error(f"PDF processing failed: {e}")
            
            # 에러 발생시에도 실시간 설정에 따라 처리
            current_config = self.config_manager.get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            if provider == 'no_model':
                try:
                    return await self._extract_text_from_pdf_fallback(file_path)
                except:
                    return "[PDF 파일: 텍스트 추출 중 오류 발생]"
            else:
                return await self._extract_text_from_pdf_via_ocr(file_path)

    async def _extract_text_from_pdf_fallback(self, file_path: str) -> str:
        """PDF 파일에서 텍스트 추출 (PyPDF2 fallback)"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n=== 페이지 {page_num + 1} ===\n"
                        text += page_text + "\n"
            return TextUtils.clean_text(text)
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            raise
    
    async def _extract_text_from_pdf_via_ocr(self, file_path: str) -> str:
        """PDF를 이미지로 변환 후 배치 OCR 메서드 사용"""
        try:
            current_config = self.config_manager.get_current_image_text_config()
            if not self.config_manager.is_image_text_enabled(current_config):
                logger.warning("OCR is disabled, falling back to text extraction")
                return await self._extract_text_from_pdf_fallback(file_path)
            
            logger.info(f"Converting PDF to images for OCR: {file_path}")
            
            image_files = await self.ocr_processor.convert_pdf_to_images(file_path)
            
            if not image_files:
                logger.warning("Failed to convert PDF to images, falling back to text extraction")
                return await self._extract_text_from_pdf_fallback(file_path)
            
            try:
                batch_size = current_config.get('batch_size', 1)
                logger.info(f"Processing {len(image_files)} pages with batch size: {batch_size}")
                
                page_texts = await self.ocr_processor.convert_images_to_text_batch(image_files, batch_size)
                
                # 결과 조합
                all_text = ""
                for i, page_text in enumerate(page_texts):
                    if not page_text.startswith("[이미지 파일:"):
                        all_text += f"\n=== 페이지 {i+1} (OCR) ===\n"
                        all_text += page_text + "\n"
                    else:
                        logger.warning(f"OCR failed for page {i+1}: {page_text}")
                       
            finally:
                # 임시 파일 정리
                for temp_file in image_files:
                    try:
                        import os
                        os.unlink(temp_file)
                    except:
                        pass
            
            if all_text.strip():
                logger.info(f"Successfully extracted text via batch OCR: {len(all_text)} chars")
                return TextUtils.clean_text(all_text)
            else:
                logger.warning("OCR failed, falling back to text extraction")
                return await self._extract_text_from_pdf_fallback(file_path)
                
        except Exception as e:
            logger.error(f"OCR processing failed for {file_path}: {e}")
            logger.warning("OCR failed, falling back to text extraction")
            return await self._extract_text_from_pdf_fallback(file_path)

    # DOCX 관련 메서드들
    async def extract_text_from_docx(self, file_path: str) -> str:
        """DOCX 파일에서 텍스트 추출 (실시간 설정에 따라 텍스트 추출 또는 OCR)"""
        try:
            current_config = self.config_manager.get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            logger.info(f"🔄 Real-time DOCX processing with provider: {provider}")
            
            if provider == 'no_model':
                logger.info("Using DOCX text extraction mode (no_model)")
                return await self._extract_text_from_docx_fallback(file_path)
            
            else:
                logger.info(f"Using DOCX OCR mode with provider: {provider}")
                return await self._extract_text_from_docx_via_ocr(file_path)
                
        except Exception as e:
            logger.error(f"DOCX processing failed: {e}")
            
            current_config = self.config_manager.get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            if provider == 'no_model':
                try:
                    return await self._extract_text_from_docx_fallback(file_path)
                except:
                    return "[DOCX 파일: 텍스트 추출 중 오류 발생]"
            else:
                try:
                    logger.warning("DOCX OCR failed, trying basic text extraction")
                    return await self._extract_text_from_docx_fallback(file_path)
                except:
                    return "[DOCX 파일: OCR 및 텍스트 추출 모두 실패]"

    async def _extract_text_from_docx_via_ocr(self, file_path: str) -> str:
        """DOCX를 이미지로 변환 후 배치 OCR 처리"""
        try:
            current_config = self.config_manager.get_current_image_text_config()
            if not self.config_manager.is_image_text_enabled(current_config):
                logger.warning("OCR is disabled for DOCX, falling back to text extraction")
                return await self._extract_text_from_docx_fallback(file_path)
            
            logger.info(f"Converting DOCX to images for OCR: {file_path}")
            
            image_files = await self.ocr_processor.convert_docx_to_images(file_path)
            
            if not image_files:
                logger.warning("Failed to convert DOCX to images, falling back to text extraction")
                return await self._extract_text_from_docx_fallback(file_path)
            
            try:
                batch_size = current_config.get('batch_size', 1)
                logger.info(f"Processing {len(image_files)} DOCX pages with batch size: {batch_size}")
                
                page_texts = await self.ocr_processor.convert_images_to_text_batch(image_files, batch_size)
                
                # 결과 조합
                all_text = ""
                for i, page_text in enumerate(page_texts):
                    if not page_text.startswith("[이미지 파일:"):
                        all_text += f"\n=== 페이지 {i+1} (OCR) ===\n"
                        all_text += page_text + "\n"
                    else:
                        logger.warning(f"OCR failed for DOCX page {i+1}: {page_text}")
                        
            finally:
                # 임시 파일 정리
                for temp_file in image_files:
                    try:
                        import os
                        os.unlink(temp_file)
                    except:
                        pass
            
            if all_text.strip():
                logger.info(f"Successfully extracted DOCX text via batch OCR: {len(all_text)} chars")
                return TextUtils.clean_text(all_text)
            else:
                logger.warning("DOCX OCR failed, falling back to text extraction")
                return await self._extract_text_from_docx_fallback(file_path)
                
        except Exception as e:
            logger.error(f"DOCX OCR processing failed for {file_path}: {e}")
            logger.warning("DOCX OCR failed, falling back to text extraction")
            return await self._extract_text_from_docx_fallback(file_path)

    async def _extract_text_from_docx_fallback(self, file_path: str) -> str:
        """DOCX 파일에서 기본 텍스트 추출 (기존 방법)"""
        try:
            doc = Document(file_path)
            text = ""
            processed_tables = set()
            
            # 방법 1: 문서의 모든 요소를 순서대로 처리 (고급 방법)
            try:
                for element in doc.element.body:
                    if element.tag.endswith('p'):
                        para_text = self._extract_paragraph_text(element, doc)
                        if para_text.strip():
                            text += para_text + "\n"
                            
                    elif element.tag.endswith('tbl'):
                        table_text = self._extract_table_text(element)
                        if table_text.strip():
                            text += "\n=== 표 ===\n" + table_text + "\n=== 표 끝 ===\n\n"
                            processed_tables.add(table_text)
                
                logger.info("Successfully used advanced DOCX parsing method")
            except Exception as e:
                logger.warning(f"Advanced parsing failed, falling back to simple method: {e}")
                # Fallback: 간단한 방법으로 모든 단락 추출
                text = ""
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text += paragraph.text + "\n"
            
            # 방법 2: 모든 표를 확실히 추출 (중복 제거)
            for i, table in enumerate(doc.tables):
                table_text = self._extract_simple_table_text(table)
                if table_text.strip():
                    # 이미 처리된 표인지 확인 (간단한 비교)
                    is_duplicate = any(
                        TextUtils.is_similar_table_text(table_text, processed) 
                        for processed in processed_tables
                    )
                    
                    if not is_duplicate:
                        text += f"\n=== 표 {i+1} ===\n" + table_text + "\n=== 표 끝 ===\n\n"
                        processed_tables.add(table_text)
            
            logger.info(f"Extracted {len(processed_tables)} tables from DOCX")
            return TextUtils.clean_text(text)
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
            raise

    def _extract_paragraph_text(self, para_element, doc) -> str:
        """단락에서 텍스트와 이미지 정보 추출"""
        try:
            text = ""
            nsmap = doc.element.nsmap if hasattr(doc.element, 'nsmap') else {}
            
            for run in para_element.findall('.//w:r', nsmap):
                # 텍스트 추출
                for text_elem in run.findall('.//w:t', nsmap):
                    if text_elem.text:
                        text += text_elem.text
                
                # 이미지 정보 추출
                for drawing in run.findall('.//w:drawing', nsmap):
                    text += " [이미지] "
                    
                    # 이미지 설명 텍스트 찾기 (가능한 경우)
                    try:
                        for desc in drawing.findall('.//wp:docPr', nsmap):
                            if desc.get('descr'):
                                text += f"[설명: {desc.get('descr')}] "
                            elif desc.get('name'):
                                text += f"[이름: {desc.get('name')}] "
                    except:
                        pass
            
            return text
        except Exception as e:
            logger.debug(f"Error extracting paragraph text: {e}")
            # Fallback methods
            try:
                text_elements = para_element.findall('.//w:t', {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'})
                return ''.join([elem.text or '' for elem in text_elements])
            except:
                try:
                    return para_element.text or ""
                except:
                    return ""
    
    def _extract_table_text(self, table_element) -> str:
        """표 요소에서 텍스트 추출"""
        try:
            text = ""
            ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
            
            for row in table_element.findall('.//w:tr', ns):
                row_text = []
                for cell in row.findall('.//w:tc', ns):
                    cell_text = ""
                    for text_elem in cell.findall('.//w:t', ns):
                        if text_elem.text:
                            cell_text += text_elem.text
                    row_text.append(cell_text.strip())
                
                if any(cell.strip() for cell in row_text):
                    text += " | ".join(row_text) + "\n"
            
            return text
        except Exception as e:
            logger.debug(f"Error extracting table text: {e}")
            return ""
    
    def _extract_simple_table_text(self, table) -> str:
        """python-docx Table 객체에서 텍스트 추출"""
        try:
            text = ""
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    row_text.append(cell_text)
                
                if any(cell.strip() for cell in row_text):
                    text += " | ".join(row_text) + "\n"
            
            return text
        except Exception as e:
            logger.debug(f"Error extracting simple table text: {e}")
            return ""

    # PPT 관련 메서드들
    async def extract_text_from_ppt(self, file_path: str) -> str:
        """PPT 파일에서 텍스트 추출 (실시간 설정에 따라 텍스트 추출 또는 OCR)"""
        try:
            current_config = self.config_manager.get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            logger.info(f"🔄 Real-time PPT processing with provider: {provider}")
            
            if provider == 'no_model':
                logger.info("Using PPT text extraction mode (no_model)")
                return await self._extract_text_from_ppt_fallback(file_path)
            
            else:
                logger.info(f"Using PPT OCR mode with provider: {provider}")
                return await self._extract_text_from_ppt_via_ocr(file_path)
                
        except Exception as e:
            logger.error(f"PPT processing failed: {e}")
            
            current_config = self.config_manager.get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            if provider == 'no_model':
                try:
                    return await self._extract_text_from_ppt_fallback(file_path)
                except:
                    return "[PPT 파일: 텍스트 추출 중 오류 발생]"
            else:
                try:
                    logger.warning("PPT OCR failed, trying basic text extraction")
                    return await self._extract_text_from_ppt_fallback(file_path)
                except:
                    return "[PPT 파일: OCR 및 텍스트 추출 모두 실패]"

    async def _extract_text_from_ppt_via_ocr(self, file_path: str) -> str:
        """PPT를 이미지로 변환 후 배치 OCR 처리"""
        try:
            current_config = self.config_manager.get_current_image_text_config()
            if not self.config_manager.is_image_text_enabled(current_config):
                logger.warning("OCR is disabled for PPT, falling back to text extraction")
                return await self._extract_text_from_ppt_fallback(file_path)
            
            logger.info(f"Converting PPT to images for OCR: {file_path}")
            
            image_files = await self.ocr_processor.convert_ppt_to_images(file_path)
            
            if not image_files:
                logger.warning("Failed to convert PPT to images, falling back to text extraction")
                return await self._extract_text_from_ppt_fallback(file_path)
            
            try:
                batch_size = current_config.get('batch_size', 1)
                logger.info(f"Processing {len(image_files)} PPT slides with batch size: {batch_size}")
                
                slide_texts = await self.ocr_processor.convert_images_to_text_batch(image_files, batch_size)
                
                # 결과 조합
                all_text = ""
                for i, slide_text in enumerate(slide_texts):
                    if not slide_text.startswith("[이미지 파일:"):
                        all_text += f"\n=== 슬라이드 {i+1} (OCR) ===\n"
                        all_text += slide_text + "\n"
                    else:
                        logger.warning(f"OCR failed for PPT slide {i+1}: {slide_text}")
                        
            finally:
                # 임시 파일 정리
                for temp_file in image_files:
                    try:
                        import os
                        os.unlink(temp_file)
                    except:
                        pass
            
            if all_text.strip():
                logger.info(f"Successfully extracted PPT text via batch OCR: {len(all_text)} chars")
                return TextUtils.clean_text(all_text)
            else:
                logger.warning("PPT OCR failed, falling back to text extraction")
                return await self._extract_text_from_ppt_fallback(file_path)
                
        except Exception as e:
            logger.error(f"PPT OCR processing failed for {file_path}: {e}")
            logger.warning("PPT OCR failed, falling back to text extraction")
            return await self._extract_text_from_ppt_fallback(file_path)

    async def _extract_text_from_ppt_fallback(self, file_path: str) -> str:
        """PPT 파일에서 기본 텍스트 추출 (python-pptx 사용)"""
        if not PYTHON_PPTX_AVAILABLE:
            raise Exception("python-pptx is required for PPT file processing but is not available")
        
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                logger.info(f"Processing slide {slide_num + 1}/{len(prs.slides)}")
                
                # 슬라이드 제목 추가
                text += f"\n=== 슬라이드 {slide_num + 1} ===\n"
                
                # 슬라이드의 모든 도형에서 텍스트 추출
                slide_content = ""
                tables_found = 0
                
                for shape in slide.shapes:
                    # 텍스트 도형 처리
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content += shape.text + "\n"
                    
                    # 표 처리
                    elif hasattr(shape, "table"):
                        tables_found += 1
                        slide_content += f"\n--- 표 {tables_found} ---\n"
                        
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                row_text.append(cell_text)
                            
                            if any(cell.strip() for cell in row_text):
                                slide_content += " | ".join(row_text) + "\n"
                        
                        slide_content += f"--- 표 {tables_found} 끝 ---\n\n"
                    
                    # 차트나 다른 객체의 경우 타입 정보만 추가
                    elif hasattr(shape, "chart"):
                        slide_content += "[차트 객체]\n"
                    elif hasattr(shape, "picture"):
                        slide_content += "[이미지 객체]\n"
                
                # 슬라이드 노트 추가 (있는 경우)
                if hasattr(slide, "notes_slide") and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_content += f"\n[슬라이드 노트]\n{notes_text}\n"
                
                # 빈 슬라이드가 아닌 경우만 추가
                if slide_content.strip():
                    text += slide_content + "\n"
                else:
                    text += "[빈 슬라이드]\n\n"
            
            logger.info(f"Extracted text from {len(prs.slides)} slides")
            return TextUtils.clean_text(text)
            
        except Exception as e:
            logger.error(f"Error extracting text from PPT {file_path}: {e}")
            raise

    # 기타 파일 형식 처리
    async def extract_text_from_excel(self, file_path: str) -> str:
        """Excel 파일에서 텍스트 추출 (openpyxl 및 xlrd 사용)"""
        if not (OPENPYXL_AVAILABLE and XLRD_AVAILABLE):
            raise Exception("openpyxl and xlrd are required for Excel file processing but are not available")
        
        try:
            text = ""
            
            # 파일 확장자에 따라 처리 방법 결정
            file_extension = file_path.lower().split('.')[-1]
            
            if file_extension in ['xlsx', 'xlsm']:
                # openpyxl로 .xlsx/.xlsm 파일 처리
                workbook = load_workbook(file_path, read_only=True, data_only=True)
                
                for sheet_name in workbook.sheetnames:
                    logger.info(f"Processing sheet: {sheet_name}")
                    worksheet = workbook[sheet_name]
                    
                    # 시트 이름 추가
                    text += f"\n=== 시트: {sheet_name} ===\n"
                    
                    # 데이터가 있는 행과 열 범위 확인
                    if worksheet.max_row == 1 and worksheet.max_column == 1:
                        # 빈 시트인 경우
                        text += "빈 시트\n\n"
                        continue
                    
                    # 첫 번째 행을 헤더로 처리
                    header_row = []
                    for col in range(1, worksheet.max_column + 1):
                        cell_value = worksheet.cell(row=1, column=col).value
                        if cell_value is not None:
                            header_row.append(str(cell_value))
                        else:
                            header_row.append("")
                    
                    if any(header_row):  # 헤더가 있는 경우
                        text += "컬럼: " + ", ".join(header_row) + "\n\n"
                    
                    # 데이터 행들 처리 (헤더 제외)
                    for row in range(2, worksheet.max_row + 1):
                        row_data = []
                        for col in range(1, worksheet.max_column + 1):
                            cell_value = worksheet.cell(row=row, column=col).value
                            if cell_value is not None:
                                row_data.append(str(cell_value))
                        
                        if row_data:  # 빈 행이 아닌 경우
                            row_text = " | ".join(row_data)
                            if row_text.strip():
                                text += row_text + "\n"
                    
                    text += "\n"
                
                workbook.close()
                
            elif file_extension == 'xls':
                # xlrd로 .xls 파일 처리
                workbook = xlrd.open_workbook(file_path)
                
                for sheet_index in range(workbook.nsheets):
                    worksheet = workbook.sheet_by_index(sheet_index)
                    sheet_name = workbook.sheet_names()[sheet_index]
                    
                    logger.info(f"Processing sheet: {sheet_name}")
                    
                    # 시트 이름 추가
                    text += f"\n=== 시트: {sheet_name} ===\n"
                    
                    if worksheet.nrows == 0:
                        text += "빈 시트\n\n"
                        continue
                    
                    # 첫 번째 행을 헤더로 처리
                    if worksheet.nrows > 0:
                        header_row = []
                        for col in range(worksheet.ncols):
                            try:
                                cell_value = worksheet.cell_value(0, col)
                                if cell_value:
                                    header_row.append(str(cell_value))
                                else:
                                    header_row.append("")
                            except IndexError:
                                break
                        
                        if any(header_row):  # 헤더가 있는 경우
                            text += "컬럼: " + ", ".join(header_row) + "\n\n"
                    
                    # 데이터 행들 처리 (헤더 제외)
                    for row in range(1, worksheet.nrows):
                        row_data = []
                        for col in range(worksheet.ncols):
                            try:
                                cell_value = worksheet.cell_value(row, col)
                                if cell_value:
                                    # xlrd에서 날짜/시간 처리
                                    if worksheet.cell_type(row, col) == xlrd.XL_CELL_DATE:
                                        try:
                                            date_value = xlrd.xldate_as_tuple(cell_value, workbook.datemode)
                                            if date_value:
                                                from datetime import datetime
                                                dt = datetime(*date_value)
                                                row_data.append(dt.strftime("%Y-%m-%d %H:%M:%S"))
                                            else:
                                                row_data.append(str(cell_value))
                                        except xlrd.XLDateError:
                                            row_data.append(str(cell_value))
                                    else:
                                        row_data.append(str(cell_value))
                            except IndexError:
                                break
                        
                        if row_data:  # 빈 행이 아닌 경우
                            row_text = " | ".join(row_data)
                            if row_text.strip():
                                text += row_text + "\n"
                    
                    text += "\n"
            else:
                raise Exception(f"Unsupported Excel file format: {file_extension}")
            
            return TextUtils.clean_text(text)
            
        except Exception as e:
            logger.error(f"Error extracting text from Excel {file_path}: {e}")
            raise

    async def extract_text_from_text_file(self, file_path: str, file_type: str) -> str:
        """텍스트 기반 파일에서 텍스트 추출 (다양한 인코딩 시도)"""
        from .constants import (
            CODE_TYPES, TEXT_TYPES, CONFIG_TYPES, SCRIPT_TYPES, 
            LOG_TYPES, WEB_TYPES
        )
        
        # 파일 카테고리 결정
        if file_type in CODE_TYPES:
            category = 'code'
        elif file_type in (TEXT_TYPES + CONFIG_TYPES + SCRIPT_TYPES + LOG_TYPES + WEB_TYPES):
            category = 'text'
        else:
            category = 'text'
        
        for encoding in ENCODINGS:
            try:
                async with aiofiles.open(file_path, 'r', encoding=encoding) as file:
                    text = await file.read()
                
                logger.info(f"Successfully read {file_path} with {encoding} encoding")
                
                # 파일 카테고리에 따라 다른 정리 방식 적용
                if category == 'code':
                    return TextUtils.clean_code_text(text, file_type)
                else:
                    return TextUtils.clean_text(text)
                    
            except UnicodeDecodeError:
                logger.debug(f"Failed to read {file_path} with {encoding} encoding, trying next...")
                continue
            except Exception as e:
                logger.error(f"Error reading file {file_path} with {encoding} encoding: {e}")
                continue
        
        # 모든 인코딩이 실패한 경우
        raise Exception(f"Could not read file {file_path} with any supported encoding")