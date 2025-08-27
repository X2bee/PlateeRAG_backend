# your_package/document_processor/ppt_handler.py
import logging, os, tempfile
from typing import Any, Dict, List
from pptx import Presentation

from .utils import clean_text
from .ocr import convert_images_to_text_batch

logger = logging.getLogger("document-processor")

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except Exception:
    PDF2IMAGE_AVAILABLE = False

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False

async def convert_ppt_to_images(file_path: str) -> List[str]:
    temp_files: List[str] = []
    try:
        if PDF2IMAGE_AVAILABLE:
            import subprocess
            with tempfile.TemporaryDirectory() as td:
                try:
                    subprocess.run(['libreoffice','--headless','--convert-to','pdf','--outdir', td, file_path],
                                   check=True, capture_output=True)
                    pdf_path = os.path.join(td, os.path.splitext(os.path.basename(file_path))[0] + '.pdf')
                    if os.path.exists(pdf_path):
                        images = convert_from_path(pdf_path, dpi=300)
                        for im in images:
                            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as ti:
                                im.save(ti.name, 'PNG')
                                temp_files.append(ti.name)
                        return temp_files
                except Exception as e:
                    logger.warning(f"LibreOffice PPT conversion failed: {e}")

        if PIL_AVAILABLE:
            logger.warning("Using fallback PIL text rendering for PPT (low quality)")
            prs = Presentation(file_path)
            for idx, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                if not slide_text.strip():
                    continue
                img_w, img_h = 1200, 900
                img = Image.new('RGB', (img_w, img_h), color='white')
                d = ImageDraw.Draw(img)
                try:
                    font = ImageFont.load_default()
                except:
                    font = None
                y = 50; lh = 25
                d.text((50, y), f"=== 슬라이드 {idx+1} ===", fill='black', font=font); y += lh*2
                for line in slide_text.splitlines():
                    if not line.strip(): continue
                    if len(line) > 80:
                        words = line.split()
                        cur = ""
                        for w in words:
                            if len(cur + w) < 80:
                                cur += w + " "
                            else:
                                d.text((50, y), cur.strip(), fill='black', font=font); y += lh; cur = w + " "
                        if cur.strip():
                            d.text((50, y), cur.strip(), fill='black', font=font); y += lh
                    else:
                        d.text((50, y), line, fill='black', font=font); y += lh
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tf:
                    img.save(tf.name, 'PNG')
                    temp_files.append(tf.name)
            return temp_files
        return []
    except Exception as e:
        logger.error(f"Error converting PPT to images: {e}")
        for p in temp_files:
            try: os.unlink(p)
            except: pass
        return []

async def extract_text_from_ppt_fallback(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    tables_found = 0
    for i, slide in enumerate(prs.slides):
        text += f"\n=== 슬라이드 {i+1} ===\n"
        slide_content = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content += shape.text + "\n"
            elif hasattr(shape, "table"):
                tables_found += 1
                slide_content += f"\n--- 표 {tables_found} ---\n"
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    if any(c.strip() for c in row_text):
                        slide_content += " | ".join(row_text) + "\n"
                slide_content += f"--- 표 {tables_found} 끝 ---\n\n"
            elif hasattr(shape, "chart"):
                slide_content += "[차트 객체]\n"
            elif hasattr(shape, "picture"):
                slide_content += "[이미지 객체]\n"
        if hasattr(slide, "notes_slide") and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_content += f"\n[슬라이드 노트]\n{notes_text}\n"
        text += (slide_content + "\n") if slide_content.strip() else "[빈 슬라이드]\n\n"
    logger.info(f"Extracted text from {len(prs.slides)} slides, found {tables_found} tables total")
    return clean_text(text)

async def extract_text_from_ppt_via_ocr(file_path: str, current_config: Dict[str, Any]) -> str:
    images = await convert_ppt_to_images(file_path)
    if not images:
        return await extract_text_from_ppt_fallback(file_path)
    try:
        batch_size = current_config.get('batch_size', 1)
        slides = await convert_images_to_text_batch(images, current_config, batch_size)
        all_text = ""
        for i, t in enumerate(slides):
            if not str(t).startswith("[이미지 파일:"):
                all_text += f"\n=== 슬라이드 {i+1} (OCR) ===\n{t}\n"
        return clean_text(all_text) if all_text.strip() else await extract_text_from_ppt_fallback(file_path)
    finally:
        for p in images:
            try: os.unlink(p)
            except: pass

async def extract_text_from_ppt(file_path: str, current_config: Dict[str, Any]) -> str:
    provider = current_config.get('provider', 'no_model')
    logger.info(f"🔄 Real-time PPT processing with provider: {provider}")
    if provider == 'no_model':
        return await extract_text_from_ppt_fallback(file_path)
    return await extract_text_from_ppt_via_ocr(file_path, current_config)
