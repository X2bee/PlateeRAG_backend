"""
OCR 처리 모듈
"""

import base64
import logging
import tempfile
import os
import re
from pathlib import Path
from typing import List, Dict, Any

from .dependencies import (
    LANGCHAIN_OPENAI_AVAILABLE, PDF2IMAGE_AVAILABLE, 
    DOCX2PDF_AVAILABLE, PIL_AVAILABLE, PYTHON_PPTX_AVAILABLE
)
from .constants import OCR_SINGLE_PROMPT, get_batch_ocr_prompt

# 기본값 설정: 조건부 import 시 사용되지 않을 경우에 대비
ChatOpenAI = None
HumanMessage = None
convert_from_path = None
Presentation = None
Image = None
ImageDraw = None
ImageFont = None

if LANGCHAIN_OPENAI_AVAILABLE:
    from langchain_openai import ChatOpenAI
    from langchain_core.messages import HumanMessage

if PDF2IMAGE_AVAILABLE:
    from pdf2image import convert_from_path

docx_to_pdf_convert = None

if DOCX2PDF_AVAILABLE:
    from docx2pdf import convert as docx_to_pdf_convert

if PYTHON_PPTX_AVAILABLE:
    from pptx import Presentation

if PIL_AVAILABLE:
    from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger("document-processor")

class OCRProcessor:
    """OCR 처리 클래스"""
    
    def __init__(self, config_manager):
        self.config_manager = config_manager
    
    async def convert_images_to_text_batch(self, image_paths: List[str], batch_size: int = 1) -> List[str]:
        """여러 이미지를 배치로 텍스트 변환"""
        current_config = self.config_manager.get_current_image_text_config()
        
        if not self.config_manager.is_image_text_enabled(current_config):
            return ["[이미지 파일: 이미지-텍스트 변환이 설정되지 않았습니다]" for _ in image_paths]
        
        # 배치 크기 제한 (1-10)
        batch_size = max(1, min(batch_size, 10))
        
        results = []
        total_batches = (len(image_paths) + batch_size - 1) // batch_size
        
        for i in range(0, len(image_paths), batch_size):
            batch_paths = image_paths[i:i+batch_size]
            batch_num = (i // batch_size) + 1
            
            logger.info(f"Processing batch {batch_num}/{total_batches} with {len(batch_paths)} images")
            
            if len(batch_paths) == 1:
                # 단일 이미지는 기존 방식 사용
                result = await self.convert_image_to_text(batch_paths[0])
                results.append(result)
            else:
                # 여러 이미지는 배치 처리
                batch_results = await self._convert_multiple_images_to_text(batch_paths, current_config)
                results.extend(batch_results)
        
        return results

    async def _convert_multiple_images_to_text(self, image_paths: List[str], config: Dict[str, Any]) -> List[str]:
        """여러 이미지를 한번에 OCR 처리"""
        try:
            # 모든 이미지를 base64로 인코딩
            image_contents = []
            for image_path in image_paths:
                with open(image_path, "rb") as image_file:
                    base64_image = base64.b64encode(image_file.read()).decode('utf-8')
                    image_contents.append(base64_image)
            
            provider = config.get('provider', 'openai')
            api_key = config.get('api_key', '')
            base_url = config.get('base_url', 'https://api.openai.com/v1')
            model = config.get('model', 'gpt-4-vision-preview')
            temperature = config.get('temperature', 0.7)
            
            logger.info(f'🔄 Using batch OCR with {len(image_paths)} images, provider: {provider}')
            
            # 프로바이더별 LLM 클라이언트 생성
            if provider == 'openai':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key,
                    base_url=base_url,
                    temperature=temperature
                )
            elif provider == 'vllm':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key or 'dummy',
                    base_url=base_url,
                    temperature=temperature
                )
            else:
                logger.error(f"Unsupported image-text provider: {provider}")
                return [f"[이미지 파일: 지원하지 않는 프로바이더 - {provider}]" for _ in image_paths]
            
            # 배치 OCR 프롬프트
            prompt = get_batch_ocr_prompt(len(image_paths))
            
            # 멀티 이미지 메시지 생성
            content = [{"type": "text", "text": prompt}]
            
            for i, base64_image in enumerate(image_contents):
                content.append({
                    "type": "image_url", 
                    "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}
                })
            
            message = HumanMessage(content=content)
            
            # 응답 생성
            response = await llm.ainvoke([message])
            response_text = response.content
            
            # 응답을 이미지별로 분할
            results = self._parse_batch_ocr_response(response_text, len(image_paths))
            
            logger.info(f"Successfully processed {len(image_paths)} images in batch")
            return results
            
        except Exception as e:
            logger.error(f"Error in batch OCR processing: {e}")
            # 실패시 개별 처리로 fallback
            logger.warning("Batch OCR failed, falling back to individual processing")
            results = []
            for image_path in image_paths:
                result = await self.convert_image_to_text(image_path)
                results.append(result)
            return results

    def _parse_batch_ocr_response(self, response_text: str, expected_count: int) -> List[str]:
        """배치 OCR 응답을 이미지별로 분할"""
        try:
            # "=== 이미지 N ===" 패턴으로 분할
            pattern = r'=== 이미지 (\d+) ===\s*(.*?)(?=\s*=== 이미지 \d+ ===|\s*$)'
            matches = re.findall(pattern, response_text, re.DOTALL)
            
            results = []
            
            if matches and len(matches) >= expected_count:
                # 매칭된 결과 사용
                for i in range(expected_count):
                    if i < len(matches):
                        _, content = matches[i]
                        results.append(content.strip())
                    else:
                        results.append("[이미지 분할 실패]")
            else:
                # 패턴 매칭 실패시 단순 분할
                logger.warning("Pattern matching failed, using simple split")
                parts = re.split(r'=== 이미지 \d+ ===', response_text)
                
                for i in range(expected_count):
                    if i + 1 < len(parts):
                        results.append(parts[i + 1].strip())
                    else:
                        results.append("[이미지 분할 실패]")
            
            # 결과 개수 맞추기
            while len(results) < expected_count:
                results.append("[이미지 처리 실패]")
            
            return results[:expected_count]
            
        except Exception as e:
            logger.error(f"Error parsing batch OCR response: {e}")
            # 실패시 동일한 응답을 모든 이미지에 적용
            return [response_text for _ in range(expected_count)]
    
    async def convert_image_to_text(self, image_path: str) -> str:
        """이미지를 텍스트로 변환"""
        current_config = self.config_manager.get_current_image_text_config()
        
        if not self.config_manager.is_image_text_enabled(current_config):
            return "[이미지 파일: 이미지-텍스트 변환이 설정되지 않았습니다]"
        
        try:
            # 이미지를 base64로 인코딩
            with open(image_path, "rb") as image_file:
                base64_image = base64.b64encode(image_file.read()).decode('utf-8')
            
            provider = current_config.get('provider', 'openai')
            api_key = current_config.get('api_key', '')
            base_url = current_config.get('base_url', 'https://api.openai.com/v1')
            model = current_config.get('model', 'gpt-4-vision-preview')
            temperature = current_config.get('temperature', 0.7)
            
            logger.info(f'🔄 Using real-time image-text provider: {provider}')
            logger.info(f'Model: {model}, Base URL: {base_url}')
            
            # 프로바이더별 LLM 클라이언트 생성
            if provider == 'openai':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key,
                    base_url=base_url,
                    temperature=temperature
                )
            elif provider == 'vllm':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key or 'dummy',
                    base_url=base_url,
                    temperature=temperature
                )
            else:
                logger.error(f"Unsupported image-text provider: {provider}")
                return f"[이미지 파일: 지원하지 않는 프로바이더 - {provider}]"
            
            # 이미지 메시지 생성
            message = HumanMessage(
                content=[
                    {"type": "text", "text": OCR_SINGLE_PROMPT},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            )
            
            # 응답 생성
            response = await llm.ainvoke([message])
            logger.info(f"Successfully converted image to text using {provider}: {Path(image_path).name}")
            return response.content
            
        except Exception as e:
            logger.error(f"Error converting image to text {image_path}: {e}")
            return f"[이미지 파일: 텍스트 변환 중 오류 발생 - {str(e)}]"

    # 이미지 변환 메서드들
    async def convert_pdf_to_images(self, file_path: str) -> List[str]:
        """PDF를 이미지로 변환하여 임시 파일 리스트 반환"""
        if not PDF2IMAGE_AVAILABLE:
            logger.error("pdf2image not available for OCR processing")
            return []
        
        try:
            logger.info(f"Converting PDF to images for OCR: {file_path}")
            
            # PDF를 이미지로 변환
            images = convert_from_path(file_path, dpi=300)
            
            temp_files = []
            
            # 모든 이미지를 임시 파일로 저장
            for i, image in enumerate(images):
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                    image.save(temp_file.name, 'PNG')
                    temp_files.append(temp_file.name)
            
            return temp_files
            
        except Exception as e:
            logger.error(f"Error converting PDF to images: {e}")
            return []

    async def convert_docx_to_images(self, file_path: str) -> List[str]:
        """DOCX를 이미지로 변환하여 임시 파일 리스트 반환"""
        temp_files = []
        
        try:
            # 방법 1: docx2pdf + pdf2image 사용 (가장 권장)
            if DOCX2PDF_AVAILABLE and PDF2IMAGE_AVAILABLE:
                logger.info("Converting DOCX to PDF, then to images using docx2pdf + pdf2image")
                
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
                    # DOCX를 PDF로 변환
                    docx_to_pdf_convert(file_path, temp_pdf.name)
                    
                    # PDF를 이미지로 변환
                    images = convert_from_path(temp_pdf.name, dpi=300)
                    
                    for i, image in enumerate(images):
                        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                            image.save(temp_img.name, 'PNG')
                            temp_files.append(temp_img.name)
                    
                    # 임시 PDF 파일 삭제
                    os.unlink(temp_pdf.name)
                    
                return temp_files
            
            # 방법 2: LibreOffice 커맨드라인 사용
            elif PDF2IMAGE_AVAILABLE:
                logger.info("Trying LibreOffice command-line conversion")
                import subprocess
                
                with tempfile.TemporaryDirectory() as temp_dir:
                    try:
                        # LibreOffice로 DOCX를 PDF로 변환
                        subprocess.run([
                            'libreoffice', '--headless', '--convert-to', 'pdf',
                            '--outdir', temp_dir, file_path
                        ], check=True, capture_output=True)
                        
                        # 변환된 PDF 파일 찾기
                        pdf_name = Path(file_path).stem + '.pdf'
                        pdf_path = os.path.join(temp_dir, pdf_name)
                        
                        if os.path.exists(pdf_path):
                            # PDF를 이미지로 변환
                            images = convert_from_path(pdf_path, dpi=300)
                            
                            for i, image in enumerate(images):
                                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                                    image.save(temp_img.name, 'PNG')
                                    temp_files.append(temp_img.name)
                            
                            return temp_files
                        
                    except (subprocess.CalledProcessError, FileNotFoundError) as e:
                        logger.warning(f"LibreOffice conversion failed: {e}")
            
            # 모든 방법이 실패한 경우
            logger.error("No available method to convert DOCX to images")
            return []
            
        except Exception as e:
            logger.error(f"Error converting DOCX to images: {e}")
            # 임시 파일 정리
            for temp_file in temp_files:
                try:
                    os.unlink(temp_file)
                except:
                    pass
            return []

    async def convert_ppt_to_images(self, file_path: str) -> List[str]:
        """PPT를 이미지로 변환하여 임시 파일 리스트 반환"""
        temp_files = []
        
        try:
            # 방법 1: LibreOffice 커맨드라인 사용 (가장 권장)
            if PDF2IMAGE_AVAILABLE:
                logger.info("Converting PPT to PDF, then to images using LibreOffice + pdf2image")
                import subprocess
                
                with tempfile.TemporaryDirectory() as temp_dir:
                    try:
                        # LibreOffice로 PPT를 PDF로 변환
                        subprocess.run([
                            'libreoffice', '--headless', '--convert-to', 'pdf',
                            '--outdir', temp_dir, file_path
                        ], check=True, capture_output=True)
                        
                        # 변환된 PDF 파일 찾기
                        pdf_name = Path(file_path).stem + '.pdf'
                        pdf_path = os.path.join(temp_dir, pdf_name)
                        
                        if os.path.exists(pdf_path):
                            # PDF를 이미지로 변환
                            images = convert_from_path(pdf_path, dpi=300)
                            
                            for i, image in enumerate(images):
                                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                                    image.save(temp_img.name, 'PNG')
                                    temp_files.append(temp_img.name)
                            
                            return temp_files
                        
                    except (subprocess.CalledProcessError, FileNotFoundError) as e:
                        logger.warning(f"LibreOffice PPT conversion failed: {e}")
            
            # 방법 2: python-pptx + PIL을 이용한 텍스트 렌더링 (fallback, 품질 낮음)
            if PIL_AVAILABLE and PYTHON_PPTX_AVAILABLE:
                logger.warning("Using fallback PIL text rendering for PPT (low quality)")
                return await self._render_ppt_text_to_images(file_path)
            
            # 모든 방법이 실패한 경우
            logger.error("No available method to convert PPT to images")
            return []
            
        except Exception as e:
            logger.error(f"Error converting PPT to images: {e}")
            # 임시 파일 정리
            for temp_file in temp_files:
                try:
                    os.unlink(temp_file)
                except:
                    pass
            return []

    async def _render_ppt_text_to_images(self, file_path: str) -> List[str]:
        """PPT 텍스트를 PIL로 이미지로 렌더링 (fallback 방법)"""
        try:
            # PPT에서 슬라이드별 텍스트 추출
            prs = Presentation(file_path)
            temp_files = []
            
            for slide_num, slide in enumerate(prs.slides):
                # 슬라이드에서 모든 텍스트 추출
                slide_text = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if not slide_text.strip():  # 빈 슬라이드 스킵
                    continue
                
                # 이미지 생성
                img_width, img_height = 1200, 900  # 슬라이드 비율 (4:3)
                img = Image.new('RGB', (img_width, img_height), color='white')
                draw = ImageDraw.Draw(img)
                
                try:
                    # 기본 폰트 사용
                    font = ImageFont.load_default()
                except:
                    font = None
                
                y_offset = 50
                line_height = 25
                
                # 슬라이드 제목 추가
                draw.text((50, y_offset), f"=== 슬라이드 {slide_num + 1} ===", 
                            fill='black', font=font)
                y_offset += line_height * 2
                
                # 텍스트 렌더링
                lines = slide_text.split('\n')
                for line in lines:
                    if line.strip() and y_offset < img_height - 50:
                        # 긴 줄은 여러 줄로 분할
                        if len(line) > 80:
                            words = line.split()
                            current_line = ""
                            for word in words:
                                if len(current_line + word) < 80:
                                    current_line += word + " "
                                else:
                                    if current_line.strip():
                                        draw.text((50, y_offset), current_line.strip(), 
                                                fill='black', font=font)
                                        y_offset += line_height
                                    current_line = word + " "
                            if current_line.strip():
                                draw.text((50, y_offset), current_line.strip(), 
                                        fill='black', font=font)
                                y_offset += line_height
                        else:
                            draw.text((50, y_offset), line, fill='black', font=font)
                            y_offset += line_height
                
                # 임시 파일로 저장
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                    img.save(temp_file.name, 'PNG')
                    temp_files.append(temp_file.name)
            
            return temp_files
            
        except Exception as e:
            logger.error(f"Error rendering PPT text to images: {e}")
            return []