"""
문서 처리 모듈

이 모듈은 다양한 형식의 문서에서 텍스트를 추출하고, 
텍스트를 벡터화에 적합한 청크로 분할하는 기능을 제공합니다.
"""

import logging
import re
import bisect
import asyncio
import os
import base64
from pathlib import Path
from typing import List, Dict, Optional, Any
import aiofiles
import PyPDF2
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter, Language
from service.llm.llm_service import LLMService

import os
import csv
try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

try:
    from pdfminer.high_level import extract_text
    PDFMINER_AVAILABLE = True
except ImportError:
    PDFMINER_AVAILABLE = False

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
    print("✅ pdfplumber available")
except Exception:
    PDFPLUMBER_AVAILABLE = False

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
    print("✅ pdf2image available")
except ImportError:
    PDF2IMAGE_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
    print("✅ PyMuPDF available")
except Exception:
    PYMUPDF_AVAILABLE = False

try:
    from langchain_openai import ChatOpenAI
    from langchain_core.messages import HumanMessage
    LANGCHAIN_OPENAI_AVAILABLE = True
except ImportError:
    LANGCHAIN_OPENAI_AVAILABLE = False

# DOCX to Image 변환을 위한 라이브러리 체크
try:
    from docx2pdf import convert as docx_to_pdf_convert
    DOCX2PDF_AVAILABLE = True
    print("✅ docx2pdf available")
except ImportError:
    DOCX2PDF_AVAILABLE = False

# PPT 처리를 위한 라이브러리 체크
try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
    print("✅ python-pptx available")
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

# 대안: python-docx + PIL을 이용한 방법
try:
    from PIL import Image, ImageDraw, ImageFont
    from io import BytesIO
    PIL_AVAILABLE = True
    print("✅ PIL available")
except ImportError:
    PIL_AVAILABLE = False

logger = logging.getLogger("document-processor")

LANGCHAIN_CODE_LANGUAGE_MAP = {
    'py': Language.PYTHON, 'js': Language.JS, 'ts': Language.TS,
    'java': Language.JAVA, 'cpp': Language.CPP, 'c': Language.CPP,
    'cs': Language.CSHARP, 'go': Language.GO, 'rs': Language.RUST,
    'php': Language.PHP, 'rb': Language.RUBY, 'swift': Language.SWIFT,
    'kt': Language.KOTLIN, 'scala': Language.SCALA,
    'html': Language.HTML, 'jsx': Language.JS, 'tsx': Language.TS,
}

class DocumentProcessor:
    """문서 처리를 담당하는 클래스"""
    def __init__(self, collection_config=None):
        # 🔥 PPT 형식 추가
        self.document_types = ['pdf', 'docx', 'doc', 'pptx', 'ppt']
        self.text_types = ['txt', 'md', 'markdown', 'rtf']
        self.code_types = ['py','js','ts','java','cpp','c','h','cs','go','rs',
                          'php','rb','swift','kt','scala','dart','r','sql',
                          'html','css','jsx','tsx','vue','svelte']
        self.config_types = ['json','yaml','yml','xml','toml','ini','cfg','conf','properties','env']
        self.data_types   = ['csv','tsv','xlsx','xls']
        self.script_types = ['sh','bat','ps1','zsh','fish']
        self.log_types    = ['log']
        self.web_types    = ['htm','xhtml']
        # 🔥 JPG 형식 추가
        self.image_types  = ['jpg','jpeg','png','gif','bmp','webp']
        self.supported_types = (
            self.document_types + self.text_types + self.code_types +
            self.config_types + self.data_types + self.script_types +
            self.log_types + self.web_types + self.image_types
        )
        self.collection_config = collection_config
        
        # 변경
        if not OPENPYXL_AVAILABLE and not XLRD_AVAILABLE:
            # openpyxl도 없고 xlrd도 없으면 Excel 지원 제거
            self.supported_types = [t for t in self.supported_types if t not in ['xlsx', 'xls']]
            logger.warning("openpyxl and xlrd not available. Excel processing disabled.")
        if not LANGCHAIN_OPENAI_AVAILABLE:
            self.supported_types = [t for t in self.supported_types if t not in self.image_types]
            logger.warning("langchain_openai not available. Image processing disabled.")
        
        # 🔥 PPT는 LibreOffice로도 처리 가능하므로 지원 목록에서 제거하지 않음
        # if not PYTHON_PPTX_AVAILABLE:
        #     self.supported_types = [t for t in self.supported_types if t not in ['pptx', 'ppt']]
        #     logger.warning("python-pptx not available. PPT processing disabled.")
            
        self.encodings = ['utf-8','utf-8-sig','cp949','euc-kr','latin-1','ascii']
        if not PDFMINER_AVAILABLE:
            logger.warning("pdfminer not available. Using PyPDF2 fallback.")
        if not PDF2IMAGE_AVAILABLE:
            logger.warning("pdf2image not available. OCR disabled.")
        if not DOCX2PDF_AVAILABLE and not PIL_AVAILABLE:
            logger.warning("docx2pdf and PIL not available. DOCX/PPT OCR disabled.")

    def _get_current_image_text_config(self) -> Dict[str, Any]:
        """실시간으로 현재 IMAGE_TEXT 설정 가져오기"""
        try:
            from main import app
            if hasattr(app.state, 'config_composer'):
                collection_config = app.state.config_composer.get_config_by_category_name("collection")
                
                # 🔥 get_env_value() 대신 직접 .value 접근
                if hasattr(collection_config, 'IMAGE_TEXT_MODEL_PROVIDER'):
                    provider_obj = getattr(collection_config, 'IMAGE_TEXT_MODEL_PROVIDER')
                    base_url_obj = getattr(collection_config, 'IMAGE_TEXT_BASE_URL')
                    api_key_obj = getattr(collection_config, 'IMAGE_TEXT_API_KEY')
                    model_obj = getattr(collection_config, 'IMAGE_TEXT_MODEL_NAME')
                    temp_obj = getattr(collection_config, 'IMAGE_TEXT_TEMPERATURE')
                    # 🔥 배치 처리 설정 추가
                    batch_size_obj = getattr(collection_config, 'IMAGE_TEXT_BATCH_SIZE', None)
                    
                    # PersistentConfig 객체에서 실제 값 추출
                    config = {
                        'provider': str(provider_obj.value if hasattr(provider_obj, 'value') else provider_obj).lower(),
                        'base_url': str(base_url_obj.value if hasattr(base_url_obj, 'value') else base_url_obj),
                        'api_key': str(api_key_obj.value if hasattr(api_key_obj, 'value') else api_key_obj),
                        'model': str(model_obj.value if hasattr(model_obj, 'value') else model_obj),
                        'temperature': float(temp_obj.value if hasattr(temp_obj, 'value') else temp_obj),
                        # 🔥 배치 크기 설정 (기본값: 1, 최대: 5)
                        'batch_size': int(batch_size_obj.value if batch_size_obj and hasattr(batch_size_obj, 'value') else 1)
                    }
                    
                    logger.info(f"🔄 Direct value access config: {config}")
                    return config
            
        except Exception as e:
            logger.error(f"🔍 Error in _get_current_image_text_config: {e}")
            import traceback
            logger.error(f"🔍 Traceback: {traceback.format_exc()}")
        
        # fallback
        logger.warning("🔍 Using fallback config")
        return {'provider': 'no_model', 'batch_size': 1}

    def _is_image_text_enabled(self, config: Dict[str, Any]) -> bool:
        """설정에 따라 OCR이 활성화되어 있는지 확인"""
        provider = config.get('provider', 'no_model')
        if provider in ('openai', 'vllm'):
            # OCR 사용 가능한 프로바이더인지 확인
            if not LANGCHAIN_OPENAI_AVAILABLE:
                logger.warning("langchain_openai not available for OCR")
                return False
            return True
        return False

    def get_supported_types(self) -> List[str]:
        """지원하는 파일 형식 목록 반환"""
        return self.supported_types.copy()
    
    def get_file_category(self, file_type: str) -> str:
        """파일 타입의 카테고리 반환"""
        file_type = file_type.lower()
        if file_type in self.document_types:
            return 'document'
        elif file_type in self.text_types:
            return 'text'
        elif file_type in self.code_types:
            return 'code'
        elif file_type in self.config_types:
            return 'config'
        elif file_type in self.data_types:
            return 'data'
        elif file_type in self.script_types:
            return 'script'
        elif file_type in self.log_types:
            return 'log'
        elif file_type in self.web_types:
            return 'web'
        elif file_type in self.image_types:
            return 'image'
        else:
            return 'unknown'
    
    def clean_text(self, text):
        """텍스트 정리"""
        if not text:
            return ""
        # 연속된 공백을 하나로 통합
        text = re.sub(r'\s+', ' ', text)
        # 연속된 줄바꿈을 두 개로 제한
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
        return text.strip()

    def _is_text_quality_sufficient(self, text: Optional[str], min_chars: int = 500, min_word_ratio: float = 0.6) -> bool:
        """간단한 휴리스틱으로 텍스트 품질 판단

        - min_chars 미만이면 낮음
        - 텍스트 내 알파벳/한글 등 단어 문자의 비율이 낮으면(이미지 OCR 잡음 가능) 낮음
        """
        try:
            if not text:
                return False
            if len(text) < min_chars:
                return False
            # 단어 문자 비율 (한글/라틴/숫자 등)
            import re
            word_chars = re.findall(r"[\w\u1100-\u11FF\u3130-\u318F\uAC00-\uD7AF]", text)
            ratio = len(word_chars) / max(1, len(text))
            return ratio >= min_word_ratio
        except Exception:
            return False
    
    def clean_code_text(self, text: str, file_type: str) -> str:
        """코드 텍스트 정리 (코드의 구조를 보존)"""
        if not text:
            return ""
        
        # 코드 파일의 경우 들여쓰기와 줄바꿈을 보존
        # 다만 파일 끝의 과도한 공백은 제거
        text = text.rstrip()
        
        # 탭을 4개의 스페이스로 변환 (일관성을 위해)
        text = text.replace('\t', '    ')
        
        return text
    
    # 🔥 배치 OCR 처리 메서드들
    async def _convert_images_to_text_batch(self, image_paths: List[str], batch_size: int = 1) -> List[str]:
        """여러 이미지를 배치로 텍스트 변환 (실시간 설정 사용)"""
        current_config = self._get_current_image_text_config()
        
        if not self._is_image_text_enabled(current_config):
            return ["[이미지 파일: 이미지-텍스트 변환이 설정되지 않았습니다]" for _ in image_paths]
        
        # 배치 크기 제한 (1-5)
        batch_size = max(1, min(batch_size, 10))
        
        results = []
        total_batches = (len(image_paths) + batch_size - 1) // batch_size
        
        for i in range(0, len(image_paths), batch_size):
            batch_paths = image_paths[i:i+batch_size]
            batch_num = (i // batch_size) + 1
            
            logger.info(f"Processing batch {batch_num}/{total_batches} with {len(batch_paths)} images")
            
            if len(batch_paths) == 1:
                # 단일 이미지는 기존 방식 사용
                result = await self._convert_image_to_text(batch_paths[0])
                results.append(result)
            else:
                # 여러 이미지는 배치 처리
                batch_results = await self._convert_multiple_images_to_text(batch_paths, current_config)
                results.extend(batch_results)
        
        return results

    async def _convert_images_to_text_batch_with_reference(
        self, 
        image_paths: List[str], 
        reference_texts: List[str], 
        batch_size: int = 1
    ) -> List[str]:
        """여러 이미지를 배치로 텍스트 변환 (기계적 파싱 텍스트 참고)"""
        current_config = self._get_current_image_text_config()
        
        if not self._is_image_text_enabled(current_config):
            return ["[이미지 파일: 이미지-텍스트 변환이 설정되지 않았습니다]" for _ in image_paths]
        
        # 배치 크기 제한 (1-5)
        batch_size = max(1, min(batch_size, 10))
        
        results = []
        total_batches = (len(image_paths) + batch_size - 1) // batch_size
        
        for i in range(0, len(image_paths), batch_size):
            batch_paths = image_paths[i:i+batch_size]
            batch_references = reference_texts[i:i+batch_size] if i+batch_size <= len(reference_texts) else reference_texts[i:]
            batch_num = (i // batch_size) + 1
            
            logger.info(f"Processing batch {batch_num}/{total_batches} with {len(batch_paths)} images (with reference)")
            
            if len(batch_paths) == 1:
                # 단일 이미지는 기존 방식 사용
                reference_text = batch_references[0] if batch_references else ""
                result = await self._convert_image_to_text_with_reference(batch_paths[0], reference_text)
                results.append(result)
            else:
                # 여러 이미지는 배치 처리
                batch_results = await self._convert_multiple_images_to_text_with_reference(
                    batch_paths, batch_references, current_config
                )
                results.extend(batch_results)
        
        return results

    async def _convert_multiple_images_to_text(self, image_paths: List[str], config: Dict[str, Any]) -> List[str]:
        """여러 이미지를 한번에 OCR 처리"""
        try:
            # 모든 이미지를 base64로 인코딩
            image_contents = []
            for image_path in image_paths:
                with open(image_path, "rb") as image_file:
                    base64_image = base64.b64encode(image_file.read()).decode('utf-8')
                    image_contents.append(base64_image)
            
            provider = config.get('provider', 'openai')
            api_key = config.get('api_key', '')
            base_url = config.get('base_url', 'https://api.openai.com/v1')
            model = config.get('model', 'gpt-4-vision-preview')
            temperature = config.get('temperature', 0.7)
            
            logger.info(f'🔄 Using batch OCR with {len(image_paths)} images, provider: {provider}')
            
            # 프로바이더별 LLM 클라이언트 생성
            if provider == 'openai':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key,
                    base_url=base_url,
                    temperature=temperature
                )
            elif provider == 'vllm':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key or 'dummy',
                    base_url=base_url,
                    temperature=temperature
                )
            else:
                logger.error(f"Unsupported image-text provider: {provider}")
                return [f"[이미지 파일: 지원하지 않는 프로바이더 - {provider}]" for _ in image_paths]
            
            # 배치 OCR 프롬프트
            prompt = f"""다음 {len(image_paths)}개의 이미지를 각각 정확한 텍스트로 변환해주세요. 

                **🔥 중요: 연속된 페이지의 표 처리 규칙**
                   이 이미지들은 연속된 문서 페이지입니다. 표가 여러 페이지에 걸쳐 나뉘어져 있을 수 있으므로 다음 규칙을 반드시 지켜주세요:
                   
                   1. **표 연속성 인식**: 이전 페이지에서 시작된 표가 다음 페이지에서 계속되는 경우를 인식해주세요
                   2. **시기 정보 보존**: 담보취득 시기(예: "2018.09.18일부터 담보취득분")가 이전 페이지에 있고 다음 페이지에 해당 데이터가 있는 경우, 시기 정보를 각 행에 포함해주세요
                   3. **병합된 셀 완전 복원**: 병합된 셀의 내용이 여러 행에 적용되는 경우, 모든 해당 행에 완전한 정보를 표시해주세요
                   4. **지역명 완전 표기**: 지역명이 생략되거나 빈 칸으로 표시된 경우, 문맥을 고려하여 완전한 지역명을 추가해주세요

                   **중요한 규칙:**
                   1. 각 이미지의 결과를 명확히 구분해주세요
                   2. 다음 형식으로 응답해주세요:

                   === 이미지 1 ===
                   [첫 번째 이미지의 텍스트 내용]

                   === 이미지 2 ===
                   [두 번째 이미지의 텍스트 내용]

                   === 이미지 3 ===
                   [세 번째 이미지의 텍스트 내용]

                   **변환 규칙:**
                   1. **표 구조 보존**: 표가 있다면 정확한 행과 열 구조를 유지하고, 마크다운 표 형식으로 변환해주세요 또한 병합된 셀의 경우, 각각 해당 사항에 모두 넣어주세요.
                   2. **레이아웃 유지**: 원본의 레이아웃, 들여쓰기, 줄바꿈을 최대한 보존해주세요
                   3. **정확한 텍스트**: 모든 문자, 숫자, 기호를 정확히 인식해주세요
                   4. **구조 정보**: 제목, 부제목, 목록, 단락 구분을 명확히 표현해주세요
                   5. **특수 형식**: 날짜, 금액, 주소, 전화번호 등의 형식을 정확히 유지해주세요
                   6. **슬라이드 구조**: 슬라이드 제목, 내용, 차트/그래프 설명을 구분해주세요
                   7. **섹션 구분**: MarkDown 형식을 철저히 준수해서 섹션 구분을 철저히 나눠주세요.
                   8. **표 구분**: 각 이미지 내에서 표는 [표 구분]으로 명확히 구분하고, 표의 제목, 설명, 표 본체, 텍스트 변환을 모두 포함해주세요
                   9. **언어*** : 한국어/영어가 아닌 문자를 포함하지 않도록 주의해주세요. 한자의 경우 한국어로 넣어주세요.
                   **섹션 구분 예시:**
                   - 서로 다른 주제나 단락 사이 
                   - 표와 다른 내용 사이
                   - 각 표는 제목부터 텍스트 설명까지 하나의 [표 구분]을 이룹니다.
                   - 단 상품의 제목 과 설명 / 표의 제목 및 설명 등 같은 내용은 하나의 색션을 구성해야합니다

                   **출력 형식 예시:**
                   === 이미지 1 ===
                   # 제목
                   [섹션 구분]
                   본문 내용이 여기에...

                   [섹션 구분]
                   ## 소제목
                   다른 섹션의 내용...

                   [표 구분]
                   ### 지역별 상가 운영 현황 (2024년 기준)
                   ※ 단위: 면적(㎡), 임대료(만원), 보증금(만원)
                   
                   | 지역 | 매장명 | 면적(㎡) | 월임대료(만원) | 보증금(만원) | 업종 | 운영상태 |
                   |------|--------|----------|----------------|--------------|------|----------|
                   | 강남구 | 카페A | 45.2 | 350 | 5,000 | 카페 | 운영중 |
                   | 강남구 | 식당B | 82.5 | 520 | 8,000 | 한식 | 운영중 |
                   | 서초구 | 의류C | 38.7 | 280 | 4,500 | 의류 | 임시휴업 |
                   | 서초구 | 편의점D | 25.3 | 180 | 2,500 | 편의점 | 운영중 |
                   | 마포구 | 학원E | 95.8 | 420 | 6,000 | 교육 | 운영중 |
                   
                   **표 내용 완전 텍스트 변환**: 이 표는 지역별 상가 운영 현황을 나타내는 표로, 2024년 기준으로 작성되었으며 면적은 제곱미터, 임대료와 보증금은 만원 단위로 표시되어 있습니다. 총 5개의 상가 정보가 포함되어 있습니다. 첫 번째 상가는 강남구에 위치한 카페A로 면적은 45.2㎡이며, 월임대료는 350만원, 보증금은 5,000만원이고 카페 업종으로 현재 운영중입니다. 두 번째 상가는 같은 강남구에 위치한 식당B로 면적은 82.5㎡이며, 월임대료는 520만원, 보증금은 8,000만원이고 한식 업종으로 현재 운영중입니다. 세 번째 상가는 서초구에 위치한 의류C로 면적은 38.7㎡이며, 월임대료는 280만원, 보증금은 4,500만원이고 의류 업종으로 현재 임시휴업 상태입니다. 네 번째 상가는 같은 서초구에 위치한 편의점D로 면적은 25.3㎡이며, 월임대료는 180만원, 보증금은 2,500만원이고 편의점 업종으로 현재 운영중입니다. 다섯 번째 상가는 마포구에 위치한 학원E로 면적은 95.8㎡이며, 월임대료는 420만원, 보증금은 6,000만원이고 교육 업종으로 현재 운영중입니다. 
                   
                   표 전체를 분석하면, 지역별로는 강남구 2개, 서초구 2개, 마포구 1개 상가가 분포되어 있습니다. 면적 범위는 25.3㎡에서 95.8㎡까지이며, 평균 면적은 약 57.5㎡입니다. 월임대료는 180만원에서 520만원까지의 범위를 보이며, 평균 월임대료는 350만원입니다. 보증금은 2,500만원에서 8,000만원까지의 범위이며, 평균 보증금은 5,200만원입니다. 업종별로는 카페, 한식, 의류, 편의점, 교육으로 다양하게 구성되어 있으며, 운영상태는 4개 상가가 운영중이고 1개 상가가 임시휴업 상태입니다.

                   텍스트만 출력하고, 추가 설명은 하지 마세요."""

            # 멀티 이미지 메시지 생성
            content = [{"type": "text", "text": prompt}]
            
            for i, base64_image in enumerate(image_contents):
                content.append({
                    "type": "image_url", 
                    "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}
                })
            
            message = HumanMessage(content=content)
            # 응답 생성
            response = await llm.ainvoke([message])
            response_text = response.content
                        # 응답을 이미지별로 분할
            results = self._parse_batch_ocr_response(response_text, len(image_paths))
            
            logger.info(f"Successfully processed {len(image_paths)} images in batch")
            return results
            
        except Exception as e:
            logger.error(f"Error in batch OCR processing: {e}")
            # 실패시 개별 처리로 fallback
            logger.warning("Batch OCR failed, falling back to individual processing")
            results = []
            for image_path in image_paths:
                result = await self._convert_image_to_text(image_path)
                results.append(result)
            return results

    async def _convert_multiple_images_to_text_with_reference(
        self, 
        image_paths: List[str], 
        reference_texts: List[str], 
        config: Dict[str, Any]
    ) -> List[str]:
        """여러 이미지를 한번에 OCR 처리 (기계적 파싱 텍스트 참고)"""
        try:
            # 모든 이미지를 base64로 인코딩
            image_contents = []
            for image_path in image_paths:
               with open(image_path, "rb") as image_file:
                   base64_image = base64.b64encode(image_file.read()).decode('utf-8')
                   image_contents.append(base64_image)
           
            provider = config.get('provider', 'openai')
            api_key = config.get('api_key', '')
            base_url = config.get('base_url', 'https://api.openai.com/v1')
            model = config.get('model', 'gpt-4-vision-preview')
            temperature = config.get('temperature', 0.7)
            
            logger.info(f'🔄 Using batch OCR with {len(image_paths)} images and reference texts, provider: {provider}')
            
            # 프로바이더별 LLM 클라이언트 생성
            if provider == 'openai':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key,
                    base_url=base_url,
                    temperature=temperature
                )
            elif provider == 'vllm':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key or 'dummy',
                    base_url=base_url,
                    temperature=temperature
                )
            else:
                logger.error(f"Unsupported image-text provider: {provider}")
                return [f"[이미지 파일: 지원하지 않는 프로바이더 - {provider}]" for _ in image_paths]
            
            # 🔥 참고 텍스트 정보 포함한 배치 OCR 프롬프트
            reference_info = ""
            if reference_texts:
                reference_info = "\n**🔥 기계적 파싱 참고 텍스트:**\n"
                for i, ref_text in enumerate(reference_texts):
                    if ref_text and ref_text.strip():
                        reference_info += f"\n--- 이미지 {i+1} 참고 텍스트 ---\n{ref_text}\n"
                    else:
                        reference_info += f"\n--- 이미지 {i+1} 참고 텍스트 ---\n[기계적 파싱으로 추출된 텍스트 없음]\n"
            
            prompt = f"""다음 {len(image_paths)}개의 이미지를 각각 정확한 텍스트로 변환해주세요.

                {reference_info}
                위의 텍스트는 해당 이미지로부터 텍스트를 파싱한 결과물입니다. 참고하여 아래의 규칙을 준수하여 응답하시오.
                    **🔥 중요: 연속된 페이지의 표 처리 규칙**
                   이 이미지들은 연속된 문서 페이지입니다. 표가 여러 페이지에 걸쳐 나뉘어져 있을 수 있으므로 다음 규칙을 반드시 지켜주세요:
                   
                   1. **표 연속성 인식**: 이전 페이지에서 시작된 표가 다음 페이지에서 계속되는 경우를 인식해주세요
                   2. **시기 정보 보존**: 담보취득 시기(예: "2018.09.18일부터 담보취득분")가 이전 페이지에 있고 다음 페이지에 해당 데이터가 있는 경우, 시기 정보를 각 행에 포함해주세요
                   3. **병합된 셀 완전 복원**: 병합된 셀의 내용이 여러 행에 적용되는 경우, 모든 해당 행에 완전한 정보를 표시해주세요
                   4. **지역명 완전 표기**: 지역명이 생략되거나 빈 칸으로 표시된 경우, 문맥을 고려하여 완전한 지역명을 추가해주세요

                   **중요한 규칙:**
                   1. 각 이미지의 결과를 명확히 구분해주세요
                   2. 다음 형식으로 응답해주세요:

                   === 이미지 1 ===
                   [첫 번째 이미지의 텍스트 내용]

                   === 이미지 2 ===
                   [두 번째 이미지의 텍스트 내용]

                   === 이미지 3 ===
                   [세 번째 이미지의 텍스트 내용]

                   **변환 규칙:**
                   1. **표 구조 보존**: 표가 있다면 정확한 행과 열 구조를 유지하고, 마크다운 표 형식으로 변환해주세요 또한 병합된 셀의 경우, 각각 해당 사항에 모두 넣어주세요.
                    => 표의 머지된 부분을 정확하게 고려해서 표 아래에 모든 내용을 텍스트로 변환하여 적어주세요                   
                    2. **레이아웃 유지**: 원본의 레이아웃, 들여쓰기, 줄바꿈을 최대한 보존해주세요
                   3. **정확한 텍스트**: 모든 문자, 숫자, 기호를 정확히 인식해주세요
                   4. **구조 정보**: 제목, 부제목, 목록, 단락 구분을 명확히 표현해주세요
                   5. **특수 형식**: 날짜, 금액, 주소, 전화번호 등의 형식을 정확히 유지해주세요
                   6. **슬라이드 구조**: 슬라이드 제목, 내용, 차트/그래프 설명을 구분해주세요
                   7. **섹션 구분**: MarkDown 형식을 철저히 준수해서 섹션 구분을 철저히 나눠주세요.
                   8. **표 구분**: 각 이미지 내에서 표는 [표 구분]으로 명확히 구분하고, 표의 제목, 설명, 표 본체, 텍스트 변환을 모두 포함해주세요
                   9. **언어*** : 한국어/영어가 아닌 문자를 포함하지 않도록 주의해주세요. 한자의 경우 한국어로 넣어주세요.
                   **섹션 구분 예시:**
                   - 서로 다른 주제나 단락 사이 
                   - 표와 다른 내용 사이
                   - 각 표는 제목부터 텍스트 설명까지 하나의 [표 구분]을 이룹니다.
                   - 단 상품의 제목 과 설명 / 표의 제목 및 설명 등 같은 내용은 하나의 색션을 구성해야합니다

                   **출력 형식 예시:**
                   === 이미지 1 ===
                   # 제목
                   [섹션 구분]
                   본문 내용이 여기에...

                   [섹션 구분]
                   ## 소제목
                   다른 섹션의 내용...

                   [표 구분]
                   ### 지역별 상가 운영 현황 (2024년 기준)
                   ※ 단위: 면적(㎡), 임대료(만원), 보증금(만원)
                   
                   | 지역 | 매장명 | 면적(㎡) | 월임대료(만원) | 보증금(만원) | 업종 | 운영상태 |
                   |------|--------|----------|----------------|--------------|------|----------|
                   | 강남구 | 카페A | 45.2 | 350 | 5,000 | 카페 | 운영중 |
                   | 강남구 | 식당B | 82.5 | 520 | 8,000 | 한식 | 운영중 |
                   | 서초구 | 의류C | 38.7 | 280 | 4,500 | 의류 | 임시휴업 |
                   | 서초구 | 편의점D | 25.3 | 180 | 2,500 | 편의점 | 운영중 |
                   | 마포구 | 학원E | 95.8 | 420 | 6,000 | 교육 | 운영중 |
                   
                   **표 내용 완전 텍스트 변환**: 이 표는 지역별 상가 운영 현황을 나타내는 표로, 2024년 기준으로 작성되었으며 면적은 제곱미터, 임대료와 보증금은 만원 단위로 표시되어 있습니다. 총 5개의 상가 정보가 포함되어 있습니다. 첫 번째 상가는 강남구에 위치한 카페A로 면적은 45.2㎡이며, 월임대료는 350만원, 보증금은 5,000만원이고 카페 업종으로 현재 운영중입니다. 두 번째 상가는 같은 강남구에 위치한 식당B로 면적은 82.5㎡이며, 월임대료는 520만원, 보증금은 8,000만원이고 한식 업종으로 현재 운영중입니다. 세 번째 상가는 서초구에 위치한 의류C로 면적은 38.7㎡이며, 월임대료는 280만원, 보증금은 4,500만원이고 의류 업종으로 현재 임시휴업 상태입니다. 네 번째 상가는 같은 서초구에 위치한 편의점D로 면적은 25.3㎡이며, 월임대료는 180만원, 보증금은 2,500만원이고 편의점 업종으로 현재 운영중입니다. 다섯 번째 상가는 마포구에 위치한 학원E로 면적은 95.8㎡이며, 월임대료는 420만원, 보증금은 6,000만원이고 교육 업종으로 현재 운영중입니다. 
                   
                   표 전체를 분석하면, 지역별로는 강남구 2개, 서초구 2개, 마포구 1개 상가가 분포되어 있습니다. 면적 범위는 25.3㎡에서 95.8㎡까지이며, 평균 면적은 약 57.5㎡입니다. 월임대료는 180만원에서 520만원까지의 범위를 보이며, 평균 월임대료는 350만원입니다. 보증금은 2,500만원에서 8,000만원까지의 범위이며, 평균 보증금은 5,200만원입니다. 업종별로는 카페, 한식, 의류, 편의점, 교육으로 다양하게 구성되어 있으며, 운영상태는 4개 상가가 운영중이고 1개 상가가 임시휴업 상태입니다.

                   반드시 한국어 및 영어로 된 텍스트만 출력하고, 추가 설명은 하지 마세요."""

            # 멀티 이미지 메시지 생성
            content = [{"type": "text", "text": prompt}]
            
            for i, base64_image in enumerate(image_contents):
                content.append({
                    "type": "image_url", 
                    "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}
                })
            
            message = HumanMessage(content=content)
            # 응답 생성
            response = await llm.ainvoke([message])
            response_text = response.content
                        # 응답을 이미지별로 분할
            results = self._parse_batch_ocr_response(response_text, len(image_paths))
            
            logger.info(f"Successfully processed {len(image_paths)} images in batch with reference texts")
            return results
           
        except Exception as e:
            logger.error(f"Error in batch OCR processing with reference: {e}")
            # 실패시 개별 처리로 fallback
            logger.warning("Batch OCR with reference failed, falling back to individual processing")
            results = []
            for i, image_path in enumerate(image_paths):
                reference_text = reference_texts[i] if i < len(reference_texts) else ""
                result = await self._convert_image_to_text_with_reference(image_path, reference_text)
                results.append(result)
            return results

    def _parse_batch_ocr_response(self, response_text: str, expected_count: int) -> List[str]:
        """배치 OCR 응답을 이미지별로 분할"""
        try:
            # "=== 이미지 N ===" 패턴으로 분할
            import re
            
            # 패턴 매칭으로 각 이미지 섹션 찾기
            pattern = r'=== 이미지 (\d+) ===\s*(.*?)(?=\s*=== 이미지 \d+ ===|\s*$)'
            matches = re.findall(pattern, response_text, re.DOTALL)
            
            results = []
            
            if matches and len(matches) >= expected_count:
                # 매칭된 결과 사용
                for i in range(expected_count):
                    if i < len(matches):
                        _, content = matches[i]
                        results.append(content.strip())
                    else:
                        results.append("[이미지 분할 실패]")
            else:
                # 패턴 매칭 실패시 단순 분할
                logger.warning("Pattern matching failed, using simple split")
                parts = re.split(r'=== 이미지 \d+ ===', response_text)
                
                for i in range(expected_count):
                    if i + 1 < len(parts):
                        results.append(parts[i + 1].strip())
                    else:
                        results.append("[이미지 분할 실패]")
            
            # 결과 개수 맞추기
            while len(results) < expected_count:
                results.append("[이미지 처리 실패]")
            
            return results[:expected_count]
            
        except Exception as e:
            logger.error(f"Error parsing batch OCR response: {e}")
            # 실패시 동일한 응답을 모든 이미지에 적용
            return [response_text for _ in range(expected_count)]
    
    async def _convert_image_to_text(self, image_path: str) -> str:
        """이미지를 텍스트로 변환 (실시간 설정 사용)"""
        # 🔥 실시간 설정 가져오기
        current_config = self._get_current_image_text_config()
        
        if not self._is_image_text_enabled(current_config):
            return "[이미지 파일: 이미지-텍스트 변환이 설정되지 않았습니다]"
        
        try:
            # 이미지를 base64로 인코딩
            with open(image_path, "rb") as image_file:
                base64_image = base64.b64encode(image_file.read()).decode('utf-8')
            
            provider = current_config.get('provider', 'openai')
            api_key = current_config.get('api_key', '')
            base_url = current_config.get('base_url', 'https://api.openai.com/v1')
            model = current_config.get('model', 'gpt-4-vision-preview')
            temperature = current_config.get('temperature', 0.7)
            
            logger.info(f'🔄 Using real-time image-text provider: {provider}')
            logger.info(f'Model: {model}, Base URL: {base_url}')
            
            # 프로바이더별 LLM 클라이언트 생성
            if provider == 'openai':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key,
                    base_url=base_url,
                    temperature=temperature
                )
            elif provider == 'vllm':
                # vLLM의 경우 OpenAI 호환 API 사용
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key or 'dummy',  # vLLM은 보통 API 키가 필요없음
                    base_url=base_url,
                    temperature=temperature
                )
            else:
                logger.error(f"Unsupported image-text provider: {provider}")
                return f"[이미지 파일: 지원하지 않는 프로바이더 - {provider}]"
            
            # OCR 프롬프트
            prompt = """이 이미지를 정확한 텍스트로 변환해주세요. 다음 규칙을 철저히 지켜주세요:

                    1. **표 구조 보존**: 표가 있다면 정확한 행과 열 구조를 유지하고, 마크다운 표 형식으로 변환해주세요 또한 병합된 셀의 경우, 각각 해당 사항에 모두 넣어주세요.
                    2. **레이아웃 유지**: 원본의 레이아웃, 들여쓰기, 줄바꿈을 최대한 보존해주세요
                    3. **정확한 텍스트**: 모든 문자, 숫자, 기호를 정확히 인식해주세요
                    4. **구조 정보**: 제목, 부제목, 목록, 단락 구분을 명확히 표현해주세요
                    5. **특수 형식**: 날짜, 금액, 주소, 전화번호 등의 형식을 정확히 유지해주세요
                    6. **슬라이드 구조**: 슬라이드 제목, 내용, 차트/그래프 설명을 구분해주세요
                    7. **섹션 구분**: MarkDown 형식을 철저히 준수해서 섹션 구분을 철저히 나눠주세요.

                    8. **표 구분**: 각 이미지 내에서 표는 [표 구분]으로 명확히 구분하고, 표의 제목, 설명, 표 본체, 텍스트 변환을 모두 포함해주세요
                    9. **언어*** : 한국어/영어가 아닌 문자를 포함하지 않도록 주의해주세요. 한자의 경우 한국어로 넣어주세요.
                    **섹션 구분 예시:**
                    - 서로 다른 주제나 단락 사이 
                    - 표와 다른 내용 사이
                    - 각 표는 제목부터 텍스트 설명까지 하나의 [표 구분]을 이룹니다.
                    - 단 상품의 제목 과 설명 / 표의 제목 및 설명 등 같은 내용은 하나의 색션을 구성해야합니다

                    **출력 형식 예시:**
                    # 제목
                    [섹션 구분]
                    본문 내용이 여기에...

                    [섹션 구분]
                    ## 소제목
                    다른 섹션의 내용...

                    [표 구분]
                    ### 상가건물 지역별 임대 현황
                    출처: 부동산 통계청 (2024년 3분기)
                    ※ 단위: 면적(㎡), 임대료(만원/월), 보증금(만원)
                    
                    | 구분 | 지역명 | 평균임대료(만원) | 보증금(만원) | 면적(㎡) | 업종분류 | 계약현황 |
                    |------|--------|------------------|--------------|----------|----------|----------|
                    | 상가 | 강남구 | 450 | 8,500 | 65.2 | 음식점 | 계약완료 |
                    | 상가 | 서초구 | 320 | 6,200 | 48.7 | 의류 | 협의중 |
                    | 상가 | 마포구 | 280 | 4,800 | 52.3 | 카페 | 계약완료 |
                    
                    **표 내용 완전 텍스트 변환**: 이 표는 상가건물 지역별 임대 현황을 나타내는 표로, 부동산 통계청에서 발표한 2024년 3분기 자료입니다. 면적은 제곱미터, 임대료는 월 단위 만원, 보증금은 만원 단위로 표시되어 있습니다. 총 3개 지역의 상가 정보가 포함되어 있습니다. 첫 번째는 강남구 상가로 평균임대료가 450만원, 보증금이 8,500만원이며, 면적은 65.2㎡이고 음식점 업종으로 계약이 완료된 상태입니다. 두 번째는 서초구 상가로 평균임대료가 320만원, 보증금이 6,200만원이며, 면적은 48.7㎡이고 의류 업종으로 현재 협의중인 상태입니다. 세 번째는 마포구 상가로 평균임대료가 280만원, 보증금이 4,800만원이며, 면적은 52.3㎡이고 카페 업종으로 계약이 완료된 상태입니다.

                    반드시 한국어 및 영어로 된 텍스트만 출력하고, 추가 설명은 하지 마세요."""
            
            # 이미지 메시지 생성
            message = HumanMessage(
                content=[
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            )
            
            # 응답 생성
            response = await llm.ainvoke([message])
            logger.info(f"Successfully converted image to text using {provider}: {Path(image_path).name}")
            return response.content
            
        except Exception as e:
            logger.error(f"Error converting image to text {image_path}: {e}")
            return f"[이미지 파일: 텍스트 변환 중 오류 발생 - {str(e)}]"

    async def _convert_image_to_text_with_reference(self, image_path: str, reference_text: str = "") -> str:
        """이미지를 텍스트로 변환 (기계적 파싱 텍스트 참고)"""
        current_config = self._get_current_image_text_config()
        
        if not self._is_image_text_enabled(current_config):
            return "[이미지 파일: 이미지-텍스트 변환이 설정되지 않았습니다]"
        
        try:
            # 이미지를 base64로 인코딩
            with open(image_path, "rb") as image_file:
                base64_image = base64.b64encode(image_file.read()).decode('utf-8')
            
            provider = current_config.get('provider', 'openai')
            api_key = current_config.get('api_key', '')
            base_url = current_config.get('base_url', 'https://api.openai.com/v1')
            model = current_config.get('model', 'gpt-4-vision-preview')
            temperature = current_config.get('temperature', 0.7)
            
            logger.info(f'🔄 Using image-text provider with reference: {provider}')
            
            # 프로바이더별 LLM 클라이언트 생성
            if provider == 'openai':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key,
                    base_url=base_url,
                    temperature=temperature
                )
            elif provider == 'vllm':
                llm = ChatOpenAI(
                    model=model,
                    openai_api_key=api_key or 'dummy',
                    base_url=base_url,
                    temperature=temperature
                )
            else:
                logger.error(f"Unsupported image-text provider: {provider}")
                return f"[이미지 파일: 지원하지 않는 프로바이더 - {provider}]"
            
            # 🔥 참고 텍스트가 있는 경우와 없는 경우로 프롬프트 구분
            if reference_text and reference_text.strip():
                # 참고 텍스트 포함 프롬프트
                prompt = f"""이 이미지를 정확한 텍스트로 변환해주세요. 

                            **🔥 중요: 기계적 파싱 참고 텍스트 활용**
                            아래는 같은 페이지에서 기계적으로 추출된 텍스트입니다. 이를 참고하여 OCR 정확도를 높여주세요:
                            {reference_text}
                                **변환 규칙:**
                            1. **참고 텍스트 활용**: 위의 참고 텍스트를 활용하여 누락된 단어나 부정확한 인식을 보완해주세요
                                이 이미지를 정확한 텍스트로 변환해주세요. 다음 규칙을 철저히 지켜주세요:

                                1. **표 구조 보존**: 표가 있다면 정확한 행과 열 구조를 유지하고, 마크다운 표 형식으로 변환해주세요 또한 병합된 셀의 경우, 각각 해당 사항에 모두 넣어주세요.
                                => 표의 머지된 부분을 정확하게 고려해서 표 아래에 모든 내용을 텍스트로 변환하여 적어주세요.
                                2. **레이아웃 유지**: 원본의 레이아웃, 들여쓰기, 줄바꿈을 최대한 보존해주세요
                                3. **정확한 텍스트**: 모든 문자, 숫자, 기호를 정확히 인식해주세요
                                4. **구조 정보**: 제목, 부제목, 목록, 단락 구분을 명확히 표현해주세요
                                5. **특수 형식**: 날짜, 금액, 주소, 전화번호 등의 형식을 정확히 유지해주세요
                                6. **슬라이드 구조**: 슬라이드 제목, 내용, 차트/그래프 설명을 구분해주세요
                                7. **섹션 구분**: MarkDown 형식을 철저히 준수해서 섹션 구분을 철저히 나눠주세요.
                                8. **참고 텍스트 활용**: 위의 참고 텍스트를 활용하여 누락된 단어나 부정확한 인식을 보완해주세요
                                9. **표 구분**: 각 이미지 내에서 표는 [표 구분]으로 명확히 구분하고, 표의 제목, 설명, 표 본체, 텍스트 변환을 모두 포함해주세요
                                10. **언어*** : 한국어/영어가 아닌 문자를 포함하지 않도록 주의해주세요. 한자의 경우 한국어로 넣어주세요.
                                **섹션 구분 예시:**
                                - 서로 다른 주제나 단락 사이 
                                - 표와 다른 내용 사이
                                - 각 표는 제목부터 텍스트 설명까지 하나의 [표 구분]을 이룹니다.
                                - 단 상품의 제목 과 설명 / 표의 제목 및 설명 등 같은 내용은 하나의 색션을 구성해야합니다

                                **출력 형식 예시:**
                                # 제목
                                [섹션 구분]
                                본문 내용이 여기에...

                                [섹션 구분]
                                ## 소제목
                                다른 섹션의 내용...

                                [표 구분]
                                ### 상가건물 지역별 임대 현황
                                출처: 부동산 통계청 (2024년 3분기)
                                ※ 단위: 면적(㎡), 임대료(만원/월), 보증금(만원)
                                
                                | 구분 | 지역명 | 평균임대료(만원) | 보증금(만원) | 면적(㎡) | 업종분류 | 계약현황 |
                                |------|--------|------------------|--------------|----------|----------|----------|
                                | 상가 | 강남구 | 450 | 8,500 | 65.2 | 음식점 | 계약완료 |
                                | 상가 | 서초구 | 320 | 6,200 | 48.7 | 의류 | 협의중 |
                                | 상가 | 마포구 | 280 | 4,800 | 52.3 | 카페 | 계약완료 |
                                
                                **표 내용 완전 텍스트 변환**: 이 표는 상가건물 지역별 임대 현황을 나타내는 표로, 부동산 통계청에서 발표한 2024년 3분기 자료입니다. 면적은 제곱미터, 임대료는 월 단위 만원, 보증금은 만원 단위로 표시되어 있습니다. 총 3개 지역의 상가 정보가 포함되어 있습니다. 첫 번째는 강남구 상가로 평균임대료가 450만원, 보증금이 8,500만원이며, 면적은 65.2㎡이고 음식점 업종으로 계약이 완료된 상태입니다. 두 번째는 서초구 상가로 평균임대료가 320만원, 보증금이 6,200만원이며, 면적은 48.7㎡이고 의류 업종으로 현재 협의중인 상태입니다. 세 번째는 마포구 상가로 평균임대료가 280만원, 보증금이 4,800만원이며, 면적은 52.3㎡이고 카페 업종으로 계약이 완료된 상태입니다.

                                반드시 한국어 및 영어로 된 텍스트만 출력하고, 추가 설명은 하지 마세요."""
            else:
                # 기존 프롬프트 (참고 텍스트 없음)
                prompt = """이 이미지를 정확한 텍스트로 변환해주세요. 다음 규칙을 철저히 지켜주세요:

                    1. **표 구조 보존**: 표가 있다면 정확한 행과 열 구조를 유지하고, 마크다운 표 형식으로 변환해주세요 또한 병합된 셀의 경우, 각각 해당 사항에 모두 넣어주세요.
                    2. **레이아웃 유지**: 원본의 레이아웃, 들여쓰기, 줄바꿈을 최대한 보존해주세요
                    3. **정확한 텍스트**: 모든 문자, 숫자, 기호를 정확히 인식해주세요
                    4. **구조 정보**: 제목, 부제목, 목록, 단락 구분을 명확히 표현해주세요
                    5. **특수 형식**: 날짜, 금액, 주소, 전화번호 등의 형식을 정확히 유지해주세요
                    6. **슬라이드 구조**: 슬라이드 제목, 내용, 차트/그래프 설명을 구분해주세요
                    7. **섹션 구분**: MarkDown 형식을 철저히 준수해서 섹션 구분을 철저히 나눠주세요.

                    8. **표 구분**: 각 이미지 내에서 표는 [표 구분]으로 명확히 구분하고, 표의 제목, 설명, 표 본체, 텍스트 변환을 모두 포함해주세요
                    9. **언어*** : 한국어/영어가 아닌 문자를 포함하지 않도록 주의해주세요. 한자의 경우 한국어로 넣어주세요.
                    **섹션 구분 예시:**
                    - 서로 다른 주제나 단락 사이 
                    - 표와 다른 내용 사이
                    - 각 표는 제목부터 텍스트 설명까지 하나의 [표 구분]을 이룹니다.
                    - 단 상품의 제목 과 설명 / 표의 제목 및 설명 등 같은 내용은 하나의 색션을 구성해야합니다

                    **출력 형식 예시:**
                    # 제목
                    [섹션 구분]
                    본문 내용이 여기에...

                    [섹션 구분]
                    ## 소제목
                    다른 섹션의 내용...

                    [표 구분]
                    ### 상가건물 지역별 임대 현황
                    출처: 부동산 통계청 (2024년 3분기)
                    ※ 단위: 면적(㎡), 임대료(만원/월), 보증금(만원)
                    
                    | 구분 | 지역명 | 평균임대료(만원) | 보증금(만원) | 면적(㎡) | 업종분류 | 계약현황 |
                    |------|--------|------------------|--------------|----------|----------|----------|
                    | 상가 | 강남구 | 450 | 8,500 | 65.2 | 음식점 | 계약완료 |
                    | 상가 | 서초구 | 320 | 6,200 | 48.7 | 의류 | 협의중 |
                    | 상가 | 마포구 | 280 | 4,800 | 52.3 | 카페 | 계약완료 |
                    
                    **표 내용 완전 텍스트 변환**: 이 표는 상가건물 지역별 임대 현황을 나타내는 표로, 부동산 통계청에서 발표한 2024년 3분기 자료입니다. 면적은 제곱미터, 임대료는 월 단위 만원, 보증금은 만원 단위로 표시되어 있습니다. 총 3개 지역의 상가 정보가 포함되어 있습니다. 첫 번째는 강남구 상가로 평균임대료가 450만원, 보증금이 8,500만원이며, 면적은 65.2㎡이고 음식점 업종으로 계약이 완료된 상태입니다. 두 번째는 서초구 상가로 평균임대료가 320만원, 보증금이 6,200만원이며, 면적은 48.7㎡이고 의류 업종으로 현재 협의중인 상태입니다. 세 번째는 마포구 상가로 평균임대료가 280만원, 보증금이 4,800만원이며, 면적은 52.3㎡이고 카페 업종으로 계약이 완료된 상태입니다.

                    반드시 한국어 및 영어로 된 텍스트만 출력하고, 추가 설명은 하지 마세요."""
            
            # 이미지 메시지 생성
            message = HumanMessage(
                content=[
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            )
            
            # 응답 생성
            response = await llm.ainvoke([message])
            
            reference_status = "with reference" if reference_text and reference_text.strip() else "without reference"
            logger.info(f"Successfully converted image to text using {provider} ({reference_status}): {Path(image_path).name}")
            
            return response.content
            
        except Exception as e:
            logger.error(f"Error converting image to text {image_path}: {e}")
            return f"[이미지 파일: 텍스트 변환 중 오류 발생 - {str(e)}]"
    
    # PDF 관련 메서드들
    async def _extract_text_pages_for_reference(self, file_path: str) -> List[str]:
        """PDF에서 기계적 파싱으로 페이지별 텍스트 추출 (OCR 참고용)"""
        try:
            page_texts = []
            
            # PyMuPDF 사용 (페이지별 텍스트 추출)
            if PYMUPDF_AVAILABLE:
                import fitz
                doc = fitz.open(file_path)
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    page_text = page.get_text("text")
                    page_texts.append(page_text.strip())
                doc.close()
                logger.info(f"Extracted reference text from {len(page_texts)} pages using PyMuPDF")
                return page_texts
            
            # PyPDF2 fallback
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    page_texts.append(page_text.strip() if page_text else "")
            
            logger.info(f"Extracted reference text from {len(page_texts)} pages using PyPDF2")
            return page_texts
            
        except Exception as e:
            logger.warning(f"Failed to extract reference text: {e}")
            return []

    async def _extract_text_from_pdf(self, file_path: str) -> str:
        """PDF 파일에서 텍스트 추출 (실시간 설정에 따라 텍스트 추출 또는 OCR)"""
        try:
            # 🔥 실시간 설정 가져오기
            current_config = self._get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            logger.info(f"🔄 Real-time PDF processing with provider: {provider}")
            
            # no_model인 경우에만 기본 텍스트 추출
            if provider == 'no_model':
                logger.info("Using text extraction mode (no_model)")

                # PyMuPDF 우선 사용: 페이지와 라인 정보를 보다 정확하게 추출
                if PYMUPDF_AVAILABLE:
                    try:
                        fitz_text = self._extract_text_from_pdf_fitz(file_path)
                        if fitz_text and fitz_text.strip():
                            logger.info(f"Text extracted via PyMuPDF: {len(fitz_text)} chars")
                            # 품질 검사: 너무 짧거나 문자 비율이 낮으면 pdfplumber/pdfminer로 폴백
                            if not self._is_text_quality_sufficient(fitz_text):
                                logger.info("PyMuPDF extraction seems low quality, attempting pdfplumber/pdfminer fallbacks")
                                # 시도 1: pdfplumber (테이블/레이아웃 보존에 강함)
                                if PDFPLUMBER_AVAILABLE:
                                    try:
                                        import pdfplumber
                                        with pdfplumber.open(file_path) as pdf:
                                            pages = [p.extract_text() or "" for p in pdf.pages]
                                        pb_text = "\n".join(pages).strip()
                                        if pb_text and self._is_text_quality_sufficient(pb_text):
                                            logger.info(f"Text extracted via pdfplumber: {len(pb_text)} chars")
                                            return self.clean_text(pb_text)
                                        else:
                                            logger.info("pdfplumber result not sufficient, will try pdfminer")
                                    except Exception as e:
                                        logger.warning(f"pdfplumber extraction failed: {e}")

                                # 시도 2: pdfminer layout extraction
                                if PDFMINER_AVAILABLE:
                                    try:
                                        layout_text = await asyncio.to_thread(self._extract_text_from_pdf_layout, file_path)
                                        if layout_text and layout_text.strip() and self._is_text_quality_sufficient(layout_text):
                                            logger.info(f"Text extracted via pdfminer layout (after PyMuPDF fallback): {len(layout_text)} chars")
                                            return self.clean_text(layout_text)
                                    except Exception as e:
                                        logger.debug(f"pdfminer fallback failed: {e}")

                            # 기본적으로 PyMuPDF 결과 반환 (품질이 충분하면)
                            if self._is_text_quality_sufficient(fitz_text):
                                return fitz_text
                            else:
                                # 품질 부족이지만 다른 방법도 실패한 경우 이후 폴백으로 진행
                                logger.info("PyMuPDF result kept as last-resort; continuing with other fallbacks")
                    except Exception as e:
                        logger.warning(f"PyMuPDF extraction failed: {e}")
                
                # 1단계: pdfminer를 사용할 수 있으면 layout 기반으로 라인 단위 추출 시도
                if PDFMINER_AVAILABLE:
                    logger.info(f"Attempting pdfminer layout extraction for {file_path}")
                    try:
                        # layout 기반 라인 추출을 우선 시도
                        layout_text = await asyncio.to_thread(self._extract_text_from_pdf_layout, file_path)
                        if layout_text and layout_text.strip():
                            cleaned_text = self.clean_text(layout_text)
                            if cleaned_text.strip():
                                logger.info(f"Text extracted via pdfminer layout: {len(cleaned_text)} chars")
                                return cleaned_text

                        # layout 실패 시 기존 per-page 병렬 추출로 폴백
                        try:
                            pdf_reader = PyPDF2.PdfReader(file_path)
                            num_pages = len(pdf_reader.pages)
                        except Exception:
                            num_pages = None

                        all_text = ""
                        if num_pages and num_pages > 0:
                            max_workers = 2
                            page_texts = await self._extract_pages_pdfminer(file_path, num_pages, max_workers)
                            for i, page_text in enumerate(page_texts):
                                if page_text and page_text.strip():
                                    all_text += f"\n=== 페이지 {i+1} ===\n"
                                    all_text += page_text + "\n"

                            cleaned_text = self.clean_text(all_text)
                            if cleaned_text.strip():
                                logger.info(f"Text extracted per-page via pdfminer: {len(cleaned_text)} chars, pages: {num_pages}")
                                return cleaned_text

                        # 최후 폴백: 전체 추출
                        text = extract_text(file_path)
                        cleaned_text = self.clean_text(text)
                        if len(cleaned_text.strip()) > 100:
                            logger.info(f"Text extracted via pdfminer (full): {len(cleaned_text)} chars")
                            return cleaned_text
                    except Exception as e:
                        logger.warning(f"pdfminer extraction failed: {e}")
                        
                # 2단계: PyMuPDF를 런타임에서 우선 시도하고, 불가능하면 PyPDF2로 폴백
                try:
                    import fitz
                    logger.info("PyMuPDF detected at runtime, attempting extraction via PyMuPDF")
                    fitz_text = await asyncio.to_thread(self._extract_text_from_pdf_fitz, file_path)
                    if fitz_text and fitz_text.strip():
                        logger.info(f"Text extracted via PyMuPDF: {len(fitz_text)} chars")
                        return fitz_text
                except Exception as e:
                    logger.debug(f"PyMuPDF runtime attempt failed or not available: {e}")

                # PyMuPDF가 없거나 실패하면 기존 PyPDF2 fallback 사용
                logger.info(f"Using PyPDF2 fallback for {file_path}")
                text = await self._extract_text_from_pdf_fallback(file_path)
                logger.info(f"Text extracted via PyPDF2: {len(text)} chars")
                return text  # no_model에서는 OCR로 넘어가지 않음
            
            else:
                # openai, vllm 등 다른 프로바이더인 경우 무조건 OCR
                logger.info(f"Using OCR mode with provider: {provider}")
                return await self._extract_text_from_pdf_via_ocr(file_path)
                
        except Exception as e:
            logger.error(f"PDF processing failed: {e}")
            
            # 에러 발생시에도 실시간 설정에 따라 처리
            current_config = self._get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            if provider == 'no_model':
                # no_model인 경우 기본 fallback만 시도
                try:
                    return await self._extract_text_from_pdf_fallback(file_path)
                except:
                    return "[PDF 파일: 텍스트 추출 중 오류 발생]"
            else:
                # 다른 프로바이더인 경우 OCR 시도
                return await self._extract_text_from_pdf_via_ocr(file_path)

    def _extract_single_page_pdfminer(self, file_path: str, page_number: int) -> Optional[str]:
        """단일 페이지를 pdfminer로 추출 (동기 실행, to_thread로 래핑됨)"""
        try:
            # pdfminer.extract_text은 page_numbers로 집합을 받음
            page_text = extract_text(file_path, page_numbers={page_number})
            return page_text
        except Exception as e:
            logger.debug(f"pdfminer single page extraction failed for page {page_number}: {e}")
            return None

    def _extract_text_from_pdf_fitz(self, file_path: str) -> str:
        """PyMuPDF(fitz)를 사용해 페이지별 텍스트를 추출하고 페이지 마커를 포함한 문자열 반환"""
        try:
            import fitz
            doc = fitz.open(file_path)
            all_text = ""
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                # get_text("text")는 페이지의 텍스트를 라인 단위로 반환
                page_text = page.get_text("text")
                # ensure we preserve original newlines; add page marker
                all_text += f"\n=== 페이지 {page_num+1} ===\n"
                all_text += page_text
                if not page_text.endswith("\n"):
                    all_text += "\n"
            doc.close()
            return all_text
        except Exception as e:
            logger.error(f"PyMuPDF extraction error: {e}")
            raise

    async def _extract_pages_pdfminer(self, file_path: str, num_pages: int, max_workers: int = 4) -> List[Optional[str]]:
        """여러 페이지를 병렬/배치로 추출하여 리스트로 반환

        Args:
            file_path: PDF 경로
            num_pages: 전체 페이지 수
            max_workers: 병렬 작업 수 (배치 크기)
        Returns:
            페이지 인덱스 순서의 텍스트 리스트(실패한 페이지는 None)
        """
        results: List[Optional[str]] = [None] * num_pages
        try:
            batch_size = max_workers
            for start in range(0, num_pages, batch_size):
                batch = list(range(start, min(start + batch_size, num_pages)))
                tasks = [asyncio.to_thread(self._extract_single_page_pdfminer, file_path, p) for p in batch]
                try:
                    completed = await asyncio.gather(*tasks)
                except Exception as e:
                    logger.debug(f"Error during concurrent pdfminer page extraction batch {batch}: {e}")
                    completed = [None] * len(batch)

                for idx, page_text in zip(batch, completed):
                    results[idx] = page_text

            return results
        except Exception as e:
            logger.error(f"Failed to extract pages via pdfminer in batches: {e}")
            return results

    def _extract_text_from_pdf_layout(self, file_path: str) -> Optional[str]:
        """pdfminer 레이아웃 분석을 사용해 페이지별 라인 단위로 텍스트를 추출

        반환값은 페이지 마커(=== 페이지 N ===)와 각 라인을 줄바꿈으로 결합한 전체 문자열입니다.
        """
        try:
            # 지역 임포트로 pdfminer 의존성 안전하게 처리
            from pdfminer.high_level import extract_pages
            from pdfminer.layout import LTTextContainer, LTTextLine, LAParams

            all_text = ""
            laparams = LAParams()
            for page_num, page_layout in enumerate(extract_pages(file_path, laparams=laparams)):
                lines = []
                for element in page_layout:
                    if isinstance(element, LTTextContainer):
                        for obj in element:
                            # LTTextLine 또는 하위 요소에서 텍스트 획득
                            try:
                                text_line = obj.get_text()
                            except Exception:
                                continue
                            if text_line and text_line.strip():
                                lines.append(text_line.rstrip('\n'))

                if lines:
                    all_text += f"\n=== 페이지 {page_num+1} ===\n"
                    for ln in lines:
                        all_text += ln + "\n"

            return all_text if all_text.strip() else None
        except Exception as e:
            logger.debug(f"PDF layout extraction failed: {e}")
            return None

    async def _extract_text_from_pdf_fallback(self, file_path: str) -> str:
        """PDF 파일에서 텍스트 추출 (PyPDF2 fallback)"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n=== 페이지 {page_num + 1} ===\n"
                        text += page_text + "\n"
            return self.clean_text(text)
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            raise
   
    async def _extract_text_from_pdf_via_ocr(self, file_path: str) -> str:
        """PDF를 이미지로 변환 후 배치 OCR 메서드 사용 (기계적 파싱 텍스트 참고)"""
        try:
            if not PDF2IMAGE_AVAILABLE:
                logger.error("pdf2image not available for OCR processing")
                return "[PDF 파일: pdf2image 라이브러리가 필요합니다]"
            
            current_config = self._get_current_image_text_config()
            if not self._is_image_text_enabled(current_config):
                logger.warning("OCR is disabled, falling back to text extraction")
                return await self._extract_text_from_pdf_fallback(file_path)
            
            import tempfile
            import os
            
            logger.info(f"Converting PDF to images for OCR: {file_path}")
            
            # 🔥 먼저 기계적 파싱 텍스트 추출
            extracted_text_pages = await self._extract_text_pages_for_reference(file_path)
            
            # PDF를 이미지로 변환
            images = convert_from_path(file_path, dpi=300)
            
            temp_files = []
            
            try:
                # 모든 이미지를 임시 파일로 저장
                for i, image in enumerate(images):
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                        image.save(temp_file.name, 'PNG')
                        temp_files.append(temp_file.name)
                
                # 🔥 배치 처리로 OCR 수행 (참고 텍스트 포함)
                batch_size = current_config.get('batch_size', 1)
                logger.info(f"Processing {len(temp_files)} pages with batch size: {batch_size} (with reference text)")
                
                page_texts = await self._convert_images_to_text_batch_with_reference(
                    temp_files, extracted_text_pages, batch_size
                )
                
                # 결과 조합
                all_text = ""
                for i, page_text in enumerate(page_texts):
                    if not page_text.startswith("[이미지 파일:"):  # 오류 메시지가 아닌 경우
                        all_text += f"\n=== 페이지 {i+1} (OCR+참고) ===\n"
                        all_text += page_text + "\n"
                    else:
                        logger.warning(f"OCR failed for page {i+1}: {page_text}")
                        
            finally:
                # 임시 파일 정리
                for temp_file in temp_files:
                    try:
                        os.unlink(temp_file)
                    except:
                        pass
            
            if all_text.strip():
                logger.info(f"Successfully extracted text via batch OCR with reference: {len(all_text)} chars")
                return self.clean_text(all_text)
            else:
                logger.warning("OCR failed, falling back to text extraction")
                return await self._extract_text_from_pdf_fallback(file_path)
                
        except Exception as e:
            logger.error(f"OCR processing failed for {file_path}: {e}")
            logger.warning("OCR failed, falling back to text extraction")
            return await self._extract_text_from_pdf_fallback(file_path)

    # DOCX 관련 메서드들
    async def _convert_docx_to_images(self, file_path: str) -> List[str]:
        """DOCX를 이미지로 변환하여 임시 파일 리스트 반환"""
        import tempfile
        import os
        
        temp_files = []
        
        try:
            # 방법 1: docx2pdf + pdf2image 사용 (가장 권장)
            if DOCX2PDF_AVAILABLE and PDF2IMAGE_AVAILABLE:
                logger.info("Converting DOCX to PDF, then to images using docx2pdf + pdf2image")
                
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
                    # DOCX를 PDF로 변환
                    docx_to_pdf_convert(file_path, temp_pdf.name)
                    
                    # PDF를 이미지로 변환
                    images = convert_from_path(temp_pdf.name, dpi=300)
                    
                    for i, image in enumerate(images):
                        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                            image.save(temp_img.name, 'PNG')
                            temp_files.append(temp_img.name)
                    
                    # 임시 PDF 파일 삭제
                    os.unlink(temp_pdf.name)
                    
                return temp_files
            
            # 방법 2: LibreOffice 커맨드라인 사용
            elif PDF2IMAGE_AVAILABLE:
                logger.info("Trying LibreOffice command-line conversion")
                import subprocess
                
                with tempfile.TemporaryDirectory() as temp_dir:
                    try:
                        # LibreOffice로 DOCX를 PDF로 변환
                        subprocess.run([
                            'libreoffice', '--headless', '--convert-to', 'pdf',
                            '--outdir', temp_dir, file_path
                        ], check=True, capture_output=True)
                        
                        # 변환된 PDF 파일 찾기
                        pdf_name = Path(file_path).stem + '.pdf'
                        pdf_path = os.path.join(temp_dir, pdf_name)
                        
                        if os.path.exists(pdf_path):
                            # PDF를 이미지로 변환
                            images = convert_from_path(pdf_path, dpi=300)
                            
                            for i, image in enumerate(images):
                                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                                    image.save(temp_img.name, 'PNG')
                                    temp_files.append(temp_img.name)
                            
                            return temp_files
                        
                    except (subprocess.CalledProcessError, FileNotFoundError) as e:
                        logger.warning(f"LibreOffice conversion failed: {e}")
            
            # 모든 방법이 실패한 경우
            logger.error("No available method to convert DOCX to images")
            return []
            
        except Exception as e:
            logger.error(f"Error converting DOCX to images: {e}")
            # 임시 파일 정리
            for temp_file in temp_files:
                try:
                    os.unlink(temp_file)
                except:
                    pass
            return []

    async def _extract_text_from_docx_via_ocr(self, file_path: str) -> str:
        """DOCX를 이미지로 변환 후 배치 OCR 처리"""
        try:
            current_config = self._get_current_image_text_config()
            if not self._is_image_text_enabled(current_config):
                logger.warning("OCR is disabled for DOCX, falling back to text extraction")
                return await self._extract_text_from_docx_fallback(file_path)
            
            logger.info(f"Converting DOCX to images for OCR: {file_path}")
            
            image_files = await self._convert_docx_to_images(file_path)
            
            if not image_files:
                logger.warning("Failed to convert DOCX to images, falling back to text extraction")
                return await self._extract_text_from_docx_fallback(file_path)
            
            try:
                # 🔥 배치 처리로 OCR 수행
                batch_size = current_config.get('batch_size', 1)
                logger.info(f"Processing {len(image_files)} DOCX pages with batch size: {batch_size}")
                
                page_texts = await self._convert_images_to_text_batch(image_files, batch_size)
                
                # 결과 조합
                all_text = ""
                for i, page_text in enumerate(page_texts):
                    if not page_text.startswith("[이미지 파일:"):  # 오류 메시지가 아닌 경우
                        all_text += f"\n=== 페이지 {i+1} (OCR) ===\n"
                        all_text += page_text + "\n"
                    else:
                        logger.warning(f"OCR failed for DOCX page {i+1}: {page_text}")
                        
            finally:
                # 임시 파일 정리
                for temp_file in image_files:
                    try:
                        import os
                        os.unlink(temp_file)
                    except:
                        pass
            
            if all_text.strip():
                logger.info(f"Successfully extracted DOCX text via batch OCR: {len(all_text)} chars")
                return self.clean_text(all_text)
            else:
                logger.warning("DOCX OCR failed, falling back to text extraction")
                return await self._extract_text_from_docx_fallback(file_path)
                
        except Exception as e:
            logger.error(f"DOCX OCR processing failed for {file_path}: {e}")
            logger.warning("DOCX OCR failed, falling back to text extraction")
            return await self._extract_text_from_docx_fallback(file_path)

    async def _extract_text_from_docx_fallback(self, file_path: str) -> str:
        """DOCX 파일에서 기본 텍스트 추출 (기존 방법)"""
        try:
            doc = Document(file_path)
            text = ""
            processed_tables = set()  # 중복 처리 방지
            current_page = 1
            
            # 방법 1: 문서의 모든 요소를 순서대로 처리 (고급 방법)
            try:
                for element in doc.element.body:
                    if element.tag.endswith('p'):  # 단락(paragraph)
                        # 페이지 브레이크 체크
                        if self._has_page_break(element):
                            current_page += 1
                            text += f"\n=== 페이지 {current_page} ===\n"
                        
                        para_text = self._extract_paragraph_text(element, doc)
                        if para_text.strip():
                            text += para_text + "\n"
                            
                    elif element.tag.endswith('tbl'):  # 표(table)
                        table_text = self._extract_table_text(element)
                        if table_text.strip():
                            text += "\n=== 표 ===\n" + table_text + "\n=== 표 끝 ===\n\n"
                            processed_tables.add(table_text)
                
                logger.info("Successfully used advanced DOCX parsing method")
            except Exception as e:
                logger.warning(f"Advanced parsing failed, falling back to simple method: {e}")
                # Fallback: 간단한 방법으로 모든 단락 추출
                text = ""
                current_page = 1
                for paragraph in doc.paragraphs:
                    # 페이지 브레이크 체크 (간단한 방법)
                    if self._paragraph_has_page_break(paragraph):
                        current_page += 1
                        text += f"\n=== 페이지 {current_page} ===\n"
                    
                    if paragraph.text.strip():
                        text += paragraph.text + "\n"
            
            # 방법 2: 모든 표를 확실히 추출 (중복 제거)
            for i, table in enumerate(doc.tables):
                table_text = self._extract_simple_table_text(table)
                if table_text.strip():
                    # 이미 처리된 표인지 확인 (간단한 비교)
                    is_duplicate = any(
                        self._is_similar_table_text(table_text, processed) 
                        for processed in processed_tables
                    )
                    
                    if not is_duplicate:
                        text += f"\n=== 표 {i+1} ===\n" + table_text + "\n=== 표 끝 ===\n\n"
                        processed_tables.add(table_text)
            
            logger.info(f"Extracted {len(processed_tables)} tables from DOCX, detected {current_page} pages")
            
            # 첫 번째 페이지 헤더 추가 (페이지가 여러 개인 경우)
            if current_page > 1 and not text.startswith("=== 페이지 1 ==="):
                text = f"=== 페이지 1 ===\n{text}"
            
            return self.clean_text(text)
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
            raise

    async def _extract_text_from_docx(self, file_path: str) -> str:
        """DOCX 파일에서 텍스트 추출 (실시간 설정에 따라 텍스트 추출 또는 OCR)"""
        try:
            # 🔥 실시간 설정 가져오기
            current_config = self._get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            logger.info(f"🔄 Real-time DOCX processing with provider: {provider}")
            
            # no_model인 경우에만 기본 텍스트 추출
            if provider == 'no_model':
                logger.info("Using DOCX text extraction mode (no_model)")
                return await self._extract_text_from_docx_fallback(file_path)
            
            else:
                # openai, vllm 등 다른 프로바이더인 경우 OCR 시도
                logger.info(f"Using DOCX OCR mode with provider: {provider}")
                return await self._extract_text_from_docx_via_ocr(file_path)
                
        except Exception as e:
            logger.error(f"DOCX processing failed: {e}")
            
            # 에러 발생시에도 실시간 설정에 따라 처리
            current_config = self._get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            if provider == 'no_model':
                # no_model인 경우 기본 fallback만 시도
                try:
                    return await self._extract_text_from_docx_fallback(file_path)
                except:
                    return "[DOCX 파일: 텍스트 추출 중 오류 발생]"
            else:
                # 다른 프로바이더인 경우 fallback으로 기본 텍스트 추출 시도
                try:
                    logger.warning("DOCX OCR failed, trying basic text extraction")
                    return await self._extract_text_from_docx_fallback(file_path)
                except:
                    return "[DOCX 파일: OCR 및 텍스트 추출 모두 실패]"

    # PPT 관련 메서드들
    async def _convert_ppt_to_images(self, file_path: str) -> List[str]:
        """PPT를 이미지로 변환하여 임시 파일 리스트 반환"""
        import tempfile
        import os
        
        temp_files = []
        
        try:
            # 방법 1: LibreOffice 커맨드라인 사용 (가장 권장)
            if PDF2IMAGE_AVAILABLE:
                logger.info("Converting PPT to PDF, then to images using LibreOffice + pdf2image")
                import subprocess
                
                with tempfile.TemporaryDirectory() as temp_dir:
                    try:
                        # LibreOffice로 PPT를 PDF로 변환
                        subprocess.run([
                            'libreoffice', '--headless', '--convert-to', 'pdf',
                            '--outdir', temp_dir, file_path
                        ], check=True, capture_output=True)
                        
                        # 변환된 PDF 파일 찾기
                        pdf_name = Path(file_path).stem + '.pdf'
                        pdf_path = os.path.join(temp_dir, pdf_name)
                        
                        if os.path.exists(pdf_path):
                            # PDF를 이미지로 변환
                            images = convert_from_path(pdf_path, dpi=300)
                            
                            for i, image in enumerate(images):
                                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                                    image.save(temp_img.name, 'PNG')
                                    temp_files.append(temp_img.name)
                            
                            return temp_files
                        
                    except (subprocess.CalledProcessError, FileNotFoundError) as e:
                        logger.warning(f"LibreOffice PPT conversion failed: {e}")
           
            # 방법 2: python-pptx + PIL을 이용한 텍스트 렌더링 (fallback, 품질 낮음)
            if PIL_AVAILABLE and PYTHON_PPTX_AVAILABLE:
                logger.warning("Using fallback PIL text rendering for PPT (low quality)")
                return await self._render_ppt_text_to_images(file_path)
            
            # 모든 방법이 실패한 경우
            logger.error("No available method to convert PPT to images")
            return []
            
        except Exception as e:
            logger.error(f"Error converting PPT to images: {e}")
            # 임시 파일 정리
            for temp_file in temp_files:
                try:
                    os.unlink(temp_file)
                except:
                    pass
            return []

    async def _render_ppt_text_to_images(self, file_path: str) -> List[str]:
        """PPT 텍스트를 PIL로 이미지로 렌더링 (fallback 방법)"""
        import tempfile
        
        try:
            # PPT에서 슬라이드별 텍스트 추출
            prs = Presentation(file_path)
            temp_files = []
            
            for slide_num, slide in enumerate(prs.slides):
                # 슬라이드에서 모든 텍스트 추출
                slide_text = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if not slide_text.strip():  # 빈 슬라이드 스킵
                    continue
                
                # 이미지 생성
                img_width, img_height = 1200, 900  # 슬라이드 비율 (4:3)
                img = Image.new('RGB', (img_width, img_height), color='white')
                draw = ImageDraw.Draw(img)
                
                try:
                    # 기본 폰트 사용
                    font = ImageFont.load_default()
                except:
                    font = None
                
                y_offset = 50
                line_height = 25
                
                # 슬라이드 제목 추가
                draw.text((50, y_offset), f"=== 슬라이드 {slide_num + 1} ===", 
                            fill='black', font=font)
                y_offset += line_height * 2
                
                # 텍스트 렌더링
                lines = slide_text.split('\n')
                for line in lines:
                    if line.strip() and y_offset < img_height - 50:
                        # 긴 줄은 여러 줄로 분할
                        if len(line) > 80:
                            words = line.split()
                            current_line = ""
                            for word in words:
                                if len(current_line + word) < 80:
                                    current_line += word + " "
                                else:
                                    if current_line.strip():
                                        draw.text((50, y_offset), current_line.strip(), 
                                                fill='black', font=font)
                                        y_offset += line_height
                                    current_line = word + " "
                            if current_line.strip():
                                draw.text((50, y_offset), current_line.strip(), 
                                        fill='black', font=font)
                                y_offset += line_height
                        else:
                            draw.text((50, y_offset), line, fill='black', font=font)
                            y_offset += line_height
                
                # 임시 파일로 저장
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                    img.save(temp_file.name, 'PNG')
                    temp_files.append(temp_file.name)
            
            return temp_files
            
        except Exception as e:
            logger.error(f"Error rendering PPT text to images: {e}")
            return []

    async def _extract_text_from_ppt_via_ocr(self, file_path: str) -> str:
        """PPT를 이미지로 변환 후 배치 OCR 처리"""
        try:
            current_config = self._get_current_image_text_config()
            if not self._is_image_text_enabled(current_config):
                logger.warning("OCR is disabled for PPT, falling back to text extraction")
                return await self._extract_text_from_ppt_fallback(file_path)
            
            logger.info(f"Converting PPT to images for OCR: {file_path}")
            
            image_files = await self._convert_ppt_to_images(file_path)
            
            if not image_files:
                logger.warning("Failed to convert PPT to images, falling back to text extraction")
                return await self._extract_text_from_ppt_fallback(file_path)
            
            try:
                # 🔥 배치 처리로 OCR 수행
                batch_size = current_config.get('batch_size', 1)
                logger.info(f"Processing {len(image_files)} PPT slides with batch size: {batch_size}")
                
                slide_texts = await self._convert_images_to_text_batch(image_files, batch_size)
                
                # 결과 조합
                all_text = ""
                for i, slide_text in enumerate(slide_texts):
                    if not slide_text.startswith("[이미지 파일:"):  # 오류 메시지가 아닌 경우
                        all_text += f"\n=== 슬라이드 {i+1} (OCR) ===\n"
                        all_text += slide_text + "\n"
                    else:
                        logger.warning(f"OCR failed for PPT slide {i+1}: {slide_text}")
                        
            finally:
                # 임시 파일 정리
                for temp_file in image_files:
                    try:
                        import os
                        os.unlink(temp_file)
                    except:
                        pass
            
            if all_text.strip():
                logger.info(f"Successfully extracted PPT text via batch OCR: {len(all_text)} chars")
                return self.clean_text(all_text)
            else:
                logger.warning("PPT OCR failed, falling back to text extraction")
                return await self._extract_text_from_ppt_fallback(file_path)
                
        except Exception as e:
            logger.error(f"PPT OCR processing failed for {file_path}: {e}")
            logger.warning("PPT OCR failed, falling back to text extraction")
            return await self._extract_text_from_ppt_fallback(file_path)

    async def _extract_text_from_ppt_fallback(self, file_path: str) -> str:
        """PPT 파일에서 기본 텍스트 추출 (python-pptx 사용)"""
        if not PYTHON_PPTX_AVAILABLE:
            raise Exception("python-pptx is required for PPT file processing but is not available")
        
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                logger.info(f"Processing slide {slide_num + 1}/{len(prs.slides)}")
                
                # 슬라이드 제목 추가
                text += f"\n=== 슬라이드 {slide_num + 1} ===\n"
                
                # 슬라이드의 모든 도형에서 텍스트 추출
                slide_content = ""
                tables_found = 0
                
                for shape in slide.shapes:
                    # 텍스트 도형 처리
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content += shape.text + "\n"
                    
                    # 표 처리
                    elif hasattr(shape, "table"):
                        tables_found += 1
                        slide_content += f"\n--- 표 {tables_found} ---\n"
                        
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                row_text.append(cell_text)
                            
                            if any(cell.strip() for cell in row_text):  # 빈 행이 아닌 경우만
                                slide_content += " | ".join(row_text) + "\n"
                        
                        slide_content += f"--- 표 {tables_found} 끝 ---\n\n"
                    
                    # 차트나 다른 객체의 경우 타입 정보만 추가
                    elif hasattr(shape, "chart"):
                        slide_content += "[차트 객체]\n"
                    elif hasattr(shape, "picture"):
                        slide_content += "[이미지 객체]\n"
                
                # 슬라이드 노트 추가 (있는 경우)
                if hasattr(slide, "notes_slide") and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_content += f"\n[슬라이드 노트]\n{notes_text}\n"
                
                # 빈 슬라이드가 아닌 경우만 추가
                if slide_content.strip():
                    text += slide_content + "\n"
                else:
                    text += "[빈 슬라이드]\n\n"
            
            logger.info(f"Extracted text from {len(prs.slides)} slides, found {tables_found} tables total")
            return self.clean_text(text)
            
        except Exception as e:
            logger.error(f"Error extracting text from PPT {file_path}: {e}")
            raise

    async def _extract_text_from_ppt(self, file_path: str) -> str:
        """PPT 파일에서 텍스트 추출 (실시간 설정에 따라 텍스트 추출 또는 OCR)"""
        try:
            # 🔥 실시간 설정 가져오기
            current_config = self._get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            logger.info(f"🔄 Real-time PPT processing with provider: {provider}")
            
            # no_model인 경우에만 기본 텍스트 추출
            if provider == 'no_model':
                logger.info("Using PPT text extraction mode (no_model)")
                return await self._extract_text_from_ppt_fallback(file_path)
            
            else:
                # openai, vllm 등 다른 프로바이더인 경우 OCR 시도
                logger.info(f"Using PPT OCR mode with provider: {provider}")
                return await self._extract_text_from_ppt_via_ocr(file_path)
                
        except Exception as e:
            logger.error(f"PPT processing failed: {e}")
            
            # 에러 발생시에도 실시간 설정에 따라 처리
            current_config = self._get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            if provider == 'no_model':
                # no_model인 경우 기본 fallback만 시도
                try:
                    return await self._extract_text_from_ppt_fallback(file_path)
                except:
                    return "[PPT 파일: 텍스트 추출 중 오류 발생]"
            else:
                # 다른 프로바이더인 경우 fallback으로 기본 텍스트 추출 시도
                try:
                    logger.warning("PPT OCR failed, trying basic text extraction")
                    return await self._extract_text_from_ppt_fallback(file_path)
                except:
                    return "[PPT 파일: OCR 및 텍스트 추출 모두 실패]"

    # 기존 헬퍼 메서드들
    def _extract_paragraph_text(self, para_element, doc) -> str:
        """단락에서 텍스트와 이미지 정보 추출"""
        try:
            text = ""
            nsmap = doc.element.nsmap if hasattr(doc.element, 'nsmap') else {}
            
            for run in para_element.findall('.//w:r', nsmap):
                # 텍스트 추출
                for text_elem in run.findall('.//w:t', nsmap):
                    if text_elem.text:
                        text += text_elem.text
                
                # 이미지 정보 추출
                for drawing in run.findall('.//w:drawing', nsmap):
                    # 이미지가 있다는 표시 추가
                    text += " [이미지] "
                    
                    # 이미지 설명 텍스트 찾기 (가능한 경우)
                    try:
                        for desc in drawing.findall('.//wp:docPr', nsmap):
                            if desc.get('descr'):
                                text += f"[설명: {desc.get('descr')}] "
                            elif desc.get('name'):
                                text += f"[이름: {desc.get('name')}] "
                    except:
                        pass  # 이미지 메타데이터 추출 실패해도 계속
            
            return text
        except Exception as e:
            logger.debug(f"Error extracting paragraph text: {e}")
            # Fallback 1: 더 간단한 방법
            try:
                text_elements = para_element.findall('.//w:t', {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'})
                return ''.join([elem.text or '' for elem in text_elements])
            except:
                # Fallback 2: 최종 방법
                try:
                    return para_element.text or ""
                except:
                    return ""
    
    def _extract_table_text(self, table_element) -> str:
        """표 요소에서 텍스트 추출"""
        try:
            text = ""
            # namespace 정의
            ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
            
            # XML에서 직접 표 데이터 추출
            for row in table_element.findall('.//w:tr', ns):
                row_text = []
                for cell in row.findall('.//w:tc', ns):
                    cell_text = ""
                    for text_elem in cell.findall('.//w:t', ns):
                        if text_elem.text:
                            cell_text += text_elem.text
                    row_text.append(cell_text.strip())
                
                if any(cell.strip() for cell in row_text):  # 빈 행이 아닌 경우만
                    text += " | ".join(row_text) + "\n"
            
            return text
        except Exception as e:
            logger.debug(f"Error extracting table text: {e}")
            return ""
    
    def _extract_simple_table_text(self, table) -> str:
        """python-docx Table 객체에서 텍스트 추출"""
        try:
            text = ""
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    row_text.append(cell_text)
                
                if any(cell.strip() for cell in row_text):  # 빈 행이 아닌 경우만
                    text += " | ".join(row_text) + "\n"
            
            return text
        except Exception as e:
            logger.debug(f"Error extracting simple table text: {e}")
            return ""
    
    def _has_page_break(self, element) -> bool:
        """XML 요소에서 페이지 브레이크 확인"""
        try:
            # Word XML에서 페이지 브레이크 확인
            # w:br 요소의 w:type이 "page"인 경우
            nsmap = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
            
            # 페이지 브레이크 찾기
            page_breaks = element.findall('.//w:br[@w:type="page"]', nsmap)
            if page_breaks:
                return True
                
            # lastRenderedPageBreak도 확인
            last_page_breaks = element.findall('.//w:lastRenderedPageBreak', nsmap)
            if last_page_breaks:
                return True
                
            return False
        except Exception as e:
            logger.debug(f"Error checking page break: {e}")
            return False
    
    def _paragraph_has_page_break(self, paragraph) -> bool:
        """python-docx Paragraph 객체에서 페이지 브레이크 확인"""
        try:
            # paragraph의 runs을 검사하여 페이지 브레이크 찾기
            for run in paragraph.runs:
                if run.element.findall('.//w:br[@w:type="page"]', 
                                        {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}):
                    return True
                    
            # paragraph._element를 직접 검사
            if hasattr(paragraph, '_element'):
                return self._has_page_break(paragraph._element)
                
            return False
        except Exception as e:
            logger.debug(f"Error checking paragraph page break: {e}")
            return False
    
    def _extract_page_mapping(self, text: str, file_extension: str) -> List[Dict[str, Any]]:
        """텍스트에서 페이지 정보 추출
        
        Args:
            text: 원본 텍스트
            file_extension: 파일 확장자
            
        Returns:
            페이지 정보 리스트 [{"page_num": 1, "start_pos": 0, "end_pos": 100}, ...]
        """
        try:
            page_mapping = []
            
            # PDF, PPT, DOCX에서 페이지 구분자 패턴 찾기
            if file_extension in ['pdf', 'ppt', 'pptx', 'docx', 'doc']:
                import re
                
                # 페이지 구분자 패턴들
                patterns = [
                    r'=== 페이지 (\d+) ===',  # PDF PyPDF2 fallback, DOCX 기본 텍스트 추출
                    r'=== 페이지 (\d+) \(OCR\) ===',  # PDF OCR
                    r'=== 페이지 (\d+) \(OCR\+참고\) ===',  # 🔥 PDF OCR with reference
                    r'=== 슬라이드 (\d+) ===',  # PPT 기본
                    r'=== 슬라이드 (\d+) \(OCR\) ===',  # PPT OCR
                ]
                
                found_pages = False
                for pattern in patterns:
                    matches = list(re.finditer(pattern, text))
                    if matches:
                        logger.info(f"Found {len(matches)} page markers with pattern: {pattern}")
                        
                        for i, match in enumerate(matches):
                            page_num = int(match.group(1))
                            start_pos = match.end()  # 페이지 제목 다음부터
                            
                            # 다음 페이지의 시작 위치 또는 텍스트 끝
                            if i + 1 < len(matches):
                                end_pos = matches[i + 1].start()
                            else:
                                end_pos = len(text)
                            
                            page_mapping.append({
                                "page_num": page_num,
                                "start_pos": start_pos,
                                "end_pos": end_pos
                            })
                        
                        # 페이지 번호 순으로 정렬
                        page_mapping.sort(key=lambda x: x["page_num"])
                        found_pages = True
                        break
                
                if not found_pages and file_extension in ['docx', 'doc']:
                    # OCR 페이지 구분자도 없으면 텍스트 길이 기준으로 가상 페이지 생성
                    # DOCX의 경우 대략 1500자당 1페이지로 추정
                    chars_per_page = 1500
                    text_length = len(text)
                    
                    if text_length > chars_per_page:
                        estimated_pages = (text_length + chars_per_page - 1) // chars_per_page
                        logger.info(f"Creating {estimated_pages} virtual pages for DOCX based on text length ({text_length} chars)")
                        
                        for page_num in range(1, estimated_pages + 1):
                            start_pos = (page_num - 1) * chars_per_page
                            end_pos = min(page_num * chars_per_page, text_length)
                            
                            page_mapping.append({
                                "page_num": page_num,
                                "start_pos": start_pos,
                                "end_pos": end_pos
                            })
                        found_pages = True
                
                if not found_pages:
                    logger.info("No page markers found, treating as single page document")
                    page_mapping = [{"page_num": 1, "start_pos": 0, "end_pos": len(text)}]
            
            elif file_extension in ['xlsx', 'xls']:
                # Excel: 시트별로 페이지 구분
                import re
                sheet_pattern = r'=== 시트: ([^=]+) ==='
                matches = list(re.finditer(sheet_pattern, text))
                
                if matches:
                    for i, match in enumerate(matches):
                        sheet_name = match.group(1).strip()
                        start_pos = match.end()
                        
                        if i + 1 < len(matches):
                            end_pos = matches[i + 1].start()
                        else:
                            end_pos = len(text)
                        
                        page_mapping.append({
                            "page_num": i + 1,
                            "start_pos": start_pos,
                            "end_pos": end_pos,
                            "sheet_name": sheet_name
                        })
                else:
                    page_mapping = [{"page_num": 1, "start_pos": 0, "end_pos": len(text)}]
            
            else:
                # 기타 파일: 1000줄당 1페이지로 가상 페이지 생성
                lines = text.split('\n')
                lines_per_page = 1000
                
                if len(lines) > lines_per_page:
                    page_count = (len(lines) + lines_per_page - 1) // lines_per_page
                    current_pos = 0
                    
                    for page_num in range(1, page_count + 1):
                        start_line = (page_num - 1) * lines_per_page
                        end_line = min(page_num * lines_per_page, len(lines))
                        
                        page_text = '\n'.join(lines[start_line:end_line])
                        start_pos = current_pos
                        end_pos = current_pos + len(page_text)
                        
                        page_mapping.append({
                            "page_num": page_num,
                            "start_pos": start_pos,
                            "end_pos": end_pos
                        })
                        
                        current_pos = end_pos + 1  # +1 for the newline
                else:
                    page_mapping = [{"page_num": 1, "start_pos": 0, "end_pos": len(text)}]
            
            return page_mapping
            
        except Exception as e:
            logger.error(f"Failed to extract page mapping: {e}")
            return [{"page_num": 1, "start_pos": 0, "end_pos": len(text)}]
    
    def _find_chunk_position(self, chunk: str, full_text: str, start_pos: int = 0) -> int:
        """청크의 전체 텍스트 내 위치 찾기"""
        try:
            # 1차: 전체 청크 텍스트로 검색
            pos = full_text.find(chunk, start_pos)
            if pos != -1:
                return pos
            
            # 2차: 청크의 첫 줄로 검색 (10자 이상인 경우만)
            chunk_lines = chunk.strip().split('\n')
            if chunk_lines and len(chunk_lines[0]) >= 10:
                first_line = chunk_lines[0].strip()
                pos = full_text.find(first_line, start_pos)
                if pos != -1:
                    # 실제 청크 시작 위치 재탐색
                    chunk_start = full_text.find(chunk[:50] if len(chunk) > 50 else chunk, pos)
                    if chunk_start != -1:
                        return chunk_start
                    return pos
            
            # 3차: 청크의 첫 50자로 검색 (10자 이상인 경우만)
            if len(chunk.strip()) >= 10:
                chunk_start = chunk.strip()[:50]
                pos = full_text.find(chunk_start, start_pos)
                if pos != -1:
                    return pos
            
            return -1
            
        except Exception as e:
            logger.debug(f"Error finding chunk position: {e}")
            return -1
    
    def _build_line_starts(self, text: str) -> List[int]:
        """텍스트에서 각 라인의 시작 인덱스 리스트 생성"""
        try:
            starts = [0]
            for idx, ch in enumerate(text):
                if ch == '\n':
                    # 다음 문자 인덱스가 라인 시작
                    if idx + 1 < len(text):
                        starts.append(idx + 1)
            return starts
        except Exception as e:
            logger.debug(f"Error building line starts: {e}")
            return [0]

    def _build_line_offset_table(self, text: str, file_extension: str) -> List[Dict[str, int]]:
        """텍스트에서 각 라인의 글로벌 오프셋(start/end)과 페이지 정보를 포함한 테이블 생성

        반환값: [{"line_num": 1, "start": 0, "end": 10, "page": 1}, ...]
        """
        try:
            lines = text.split('\n')
            table: List[Dict[str, int]] = []
            pos = 0

            # 페이지 매핑을 미리 생성해두면 라인별로 페이지를 빠르게 찾을 수 있음
            page_mapping = self._extract_page_mapping(text, file_extension)

            def _page_for_pos(p: int) -> int:
                try:
                    for pinfo in page_mapping:
                        if pinfo["start_pos"] <= p < pinfo["end_pos"]:
                            return pinfo["page_num"]
                except Exception:
                    pass
                return 1

            for i, line in enumerate(lines):
                start = pos
                end = pos + len(line)
                mid = start + max(0, (end - start) // 2)
                page = _page_for_pos(mid)

                table.append({
                    "line_num": i + 1,
                    "start": start,
                    "end": end,
                    "page": page,
                })

                # +1 for the '\n' character that was removed by split (except maybe last line)
                pos = end + 1

            return table
        except Exception as e:
            logger.debug(f"Error building line offset table: {e}")
            return [{"line_num": 1, "start": 0, "end": len(text), "page": 1}]

    def _find_line_index_by_pos(self, pos: int, line_table: List[Dict[str, int]]) -> int:
        """이분탐색으로 위치에 해당하는 라인 인덱스(0-based)를 반환"""
        try:
            if not line_table:
                return 0
            starts = [l["start"] for l in line_table]
            idx = bisect.bisect_right(starts, pos) - 1
            if idx < 0:
                return 0
            return min(idx, len(line_table) - 1)
        except Exception as e:
            logger.debug(f"Error finding line index by pos: {e}")
            return 0

    def _build_sentence_starts(self, text: str) -> List[int]:
        """문장 단위로 텍스트를 분리해 각 문장의 시작 인덱스 리스트를 생성

        한국어의 경우 문장 종결어미(예: '다.')를 기준으로 분리하여
        문장 단위의 라인 번호를 제공함으로써 PDF에서의 시각적 라인과
        유사한 구분을 얻습니다.
        """
        try:
            starts: List[int] = []
            if not text:
                return [0]

            # 정규식: 가능한 문장 종결 패턴(다. 또는 영어권 . ? !)을 포괄
            # non-greedy로 문장 단위를 캡쳐
            pattern = re.compile(r'.*?(?:다\.|[.?!])(?=\s+|$)', re.DOTALL)
            pos = 0
            for m in pattern.finditer(text):
                start = m.start()   
                # 첫 문장의 시작이 0이 아닐 수 있으므로 보장
                if not starts or start != starts[-1]:
                    starts.append(start)
                pos = m.end()

            # 남은 텍스트가 있으면 마지막 시작 추가
            if not starts:
                return [0]
            # Ensure first is 0
            if starts[0] != 0:
                starts.insert(0, 0)
            return starts
        except Exception as e:
            logger.debug(f"Error building sentence starts: {e}")
            return [0]

    def _pos_to_line(self, pos: int, line_starts: List[int]) -> int:
        """문자열의 포지션을 라인 번호(1-based)로 변환"""
        try:
            if pos < 0:
                return 1
            # bisect로 가장 큰 start <= pos를 찾음
            idx = bisect.bisect_right(line_starts, pos) - 1
            return max(1, idx + 1)
        except Exception as e:
            logger.debug(f"Error mapping pos to line: {e}")
            return 1
    
    def _get_page_number_for_chunk(self, chunk: str, chunk_pos: int, page_mapping: List[Dict[str, Any]]) -> int:
        """청크의 페이지 번호 계산"""
        try:
            if not page_mapping:
                return 1
            
            if chunk_pos == -1:
                # 위치를 찾지 못한 경우 첫 번째 페이지로 설정
                return 1
            
            # 청크의 중간 위치를 기준으로 페이지 결정
            chunk_mid_pos = chunk_pos + len(chunk) // 2
            
            for page_info in page_mapping:
                if page_info["start_pos"] <= chunk_mid_pos < page_info["end_pos"]:
                    return page_info["page_num"]
            
            # 매핑에서 찾지 못한 경우 가장 가까운 페이지 선택
            closest_page = min(page_mapping, key=lambda p: abs(p["start_pos"] - chunk_pos))
            logger.debug(f"Chunk position {chunk_pos} not in any page range, using closest page {closest_page['page_num']}")
            return closest_page["page_num"]
            
        except Exception as e:
            logger.error(f"Error calculating page number for chunk: {e}")
            return 1
    
    async def extract_text_from_file(self, file_path: str, file_extension: str) -> str:
        """파일에서 텍스트 추출"""
        category = self.get_file_category(file_extension)
        logger.info(f"Extracting text from {file_extension} file ({category} category): {file_path}")

        if file_extension == 'pdf':
            return await self._extract_text_from_pdf(file_path)
        elif file_extension in ['docx', 'doc']:
            return await self._extract_text_from_docx(file_path)
        elif file_extension in ['pptx', 'ppt']:
            return await self._extract_text_from_ppt(file_path)
        elif file_extension in ['xlsx', 'xls']:
            return await self._extract_text_from_excel(file_path)
        elif file_extension in ['csv', 'tsv']:
            return await self._extract_text_from_csv(file_path)
        elif file_extension in self.image_types:
            return await self._convert_image_to_text(file_path)
        elif file_extension in (self.text_types + self.code_types + self.config_types + 
                                self.script_types + self.log_types + self.web_types):
            return await self._extract_text_from_text_file(file_path, file_extension)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

    def _is_similar_table_text(self, text1: str, text2: str, threshold: float = 0.8) -> bool:
        """두 표 텍스트가 유사한지 확인 (중복 제거용)"""
        try:
            if not text1 or not text2:
                return False
            
            # 공백 정규화
            text1_clean = re.sub(r'\s+', ' ', text1.strip())
            text2_clean = re.sub(r'\s+', ' ', text2.strip())
            
            # 완전히 동일한 경우
            if text1_clean == text2_clean:
                return True
            
            # 길이가 매우 다른 경우
            len1, len2 = len(text1_clean), len(text2_clean)
            if min(len1, len2) / max(len1, len2) < 0.5:
                return False
            
            # 간단한 유사도 검사 (공통 부분 비율)
            shorter = text1_clean if len1 < len2 else text2_clean
            longer = text2_clean if len1 < len2 else text1_clean
            
            common_ratio = len(set(shorter.split()) & set(longer.split())) / len(set(shorter.split()))
            return common_ratio >= threshold
            
        except Exception:
            return False
    
    async def _extract_text_from_excel(self, file_path: str) -> str:
        """Excel(xlsx/xls) 파일에서 텍스트 추출 (pandas 없이)"""
        ext = os.path.splitext(file_path)[1].lower()
        text = ""

        try:
            if ext == ".xlsx":
                if not OPENPYXL_AVAILABLE:
                    raise Exception("openpyxl이 설치되어야 .xlsx 파일을 처리할 수 있습니다.")
                wb = load_workbook(file_path, data_only=True)
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    text += f"\n=== 시트: {sheet_name} ===\n"
                    for row in ws.iter_rows(values_only=True):
                        row_text = " | ".join(str(v) for v in row if v is not None)
                        if row_text.strip():
                            text += row_text + "\n"

            elif ext == ".xls":
                if not XLRD_AVAILABLE:
                    raise Exception("xlrd가 설치되어야 .xls 파일을 처리할 수 있습니다.")
                wb = xlrd.open_workbook(file_path)
                for sheet in wb.sheets():
                    text += f"\n=== 시트: {sheet.name} ===\n"
                    for row_idx in range(sheet.nrows):
                        row_values = sheet.row_values(row_idx)
                        row_text = " | ".join(str(v) for v in row_values if v)
                        if row_text.strip():
                            text += row_text + "\n"
            else:
                raise Exception(f"지원하지 않는 Excel 형식입니다: {ext}")

            return self.clean_text(text)
        except Exception as e:
            logger.error(f"Error extracting text from Excel {file_path}: {e}")
            raise

    async def _extract_text_from_csv(self, file_path: str, encoding: str = "utf-8") -> str:
        """CSV 파일에서 텍스트 추출 (pandas 없이)"""
        text = ""
        try:
            async with aiofiles.open(file_path, mode="r", encoding=encoding) as f:
                reader = csv.reader((await f.read()).splitlines())
                for row in reader:
                    row_text = " | ".join(row)
                    if row_text.strip():
                        text += row_text + "\n"
            return self.clean_text(text)
        except Exception as e:
            logger.error(f"Error extracting text from CSV {file_path}: {e}")
            raise
    
    async def _extract_text_from_text_file(self, file_path: str, file_type: str) -> str:
        """텍스트 기반 파일에서 텍스트 추출 (다양한 인코딩 시도)"""
        category = self.get_file_category(file_type)
        
        for encoding in self.encodings:
            try:
                async with aiofiles.open(file_path, 'r', encoding=encoding) as file:
                    text = await file.read()
                
                logger.info(f"Successfully read {file_path} with {encoding} encoding")
                
                # 파일 카테고리에 따라 다른 정리 방식 적용
                if category == 'code':
                    return self.clean_code_text(text, file_type)
                else:
                    return self.clean_text(text)
                    
            except UnicodeDecodeError:
                logger.debug(f"Failed to read {file_path} with {encoding} encoding, trying next...")
                continue
            except Exception as e:
                logger.error(f"Error reading file {file_path} with {encoding} encoding: {e}")
                continue
        
        # 모든 인코딩이 실패한 경우
        raise Exception(f"Could not read file {file_path} with any supported encoding")
    
    def chunk_text(self, text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
        """텍스트를 청크로 분할
        
        Args:
            text: 분할할 텍스트
            chunk_size: 청크 크기
            chunk_overlap: 청크 간 중복 크기
            
        Returns:
            분할된 텍스트 청크 리스트
            
        Raises:
            Exception: 텍스트 분할 오류
        """
        try:
            if not text or not text.strip():
                logger.warning("Empty text provided for chunking")
                return [""]
            
            # \n\n\n이 포함되어 있으면 먼저 이것으로 분할
            # [색션 구분] 또는 [표 구분]이 포함되어 있으면 먼저 이것으로 분할
            if "[색션 구분]" in text or "[표 구분]" in text or "[섹션 구분]" in text:
                logger.info("Found section markers in text, splitting by major sections first")
                
                # 세 가지 구분자를 모두 고려해서 분할
                # 먼저 [색션 구분] 또는 [섹션 구분]으로 분할하고, 각 섹션에서 [표 구분]으로 추가 분할
                
                # 단계 1: [색션 구분] 또는 [섹션 구분]으로 먼저 분할
                if "[색션 구분]" in text:
                    temp_sections = text.split("[색션 구분]")
                elif "[섹션 구분]" in text:
                    temp_sections = text.split("[섹션 구분]")
                else:
                    temp_sections = [text]  # [표 구분]만 있는 경우
                
                major_sections = []
                
                # 단계 2: 각 섹션에서 [표 구분]으로 추가 분할
                for section in temp_sections:
                    if "[표 구분]" in section:
                        table_sections = section.split("[표 구분]")
                        major_sections.extend([s.strip() for s in table_sections if s.strip()])
                    else:
                        if section.strip():
                            major_sections.append(section.strip())
                
                logger.info(f"Split into {len(major_sections)} sections using markers")
                
                # 작은 섹션들을 합치기
                merged_sections = []
                current_merged = ""
                
                for i, section in enumerate(major_sections):
                    # 현재 합쳐진 것과 새 섹션을 합쳤을 때의 길이 계산
                    if current_merged:
                        potential_merged = current_merged + "\n\n" + section
                    else:
                        potential_merged = section
                    
                    if len(potential_merged) <= chunk_size:
                        # chunk_size를 넘지 않으면 계속 합치기
                        current_merged = potential_merged
                        logger.info(f"Merging section {i+1} (total length: {len(current_merged)})")
                    else:
                        # chunk_size를 넘으면 이전까지 합친 것을 저장하고 새로 시작
                        if current_merged:
                            merged_sections.append(current_merged)
                            logger.info(f"Added merged section with length: {len(current_merged)}")
                        current_merged = section
                
                # 마지막 섹션 추가
                if current_merged:
                    merged_sections.append(current_merged)
                    logger.info(f"Added final merged section with length: {len(current_merged)}")
                
                # 합쳐진 섹션들을 최종 청킹
                all_chunks = []
                for i, section in enumerate(merged_sections):
                    logger.info(f"Processing merged section {i+1}/{len(merged_sections)} (length: {len(section)})")
                    
                    # 섹션이 chunk_size의 2배를 초과하면 추가로 청킹
                    if len(section) > chunk_size * 2:
                        logger.info(f"Merged section {i+1} is too large ({len(section)} chars), splitting further")
                        text_splitter = RecursiveCharacterTextSplitter(
                            chunk_size=chunk_size,
                            chunk_overlap=chunk_overlap,
                            length_function=len,
                            separators=["\n\n", "\n", " ", ""]
                        )
                        section_chunks = text_splitter.split_text(section)
                        all_chunks.extend(section_chunks)
                    else:
                        # 섹션이 적당한 크기면 그대로 사용
                        all_chunks.append(section)
                
                logger.info(f"Text split into {len(all_chunks)} chunks after merging small sections")
                return all_chunks

            else:
                # 기존 방식으로 청킹
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap,
                    length_function=len,
                    separators=["\n\n", "\n", " ", ""]
                )
                chunks = text_splitter.split_text(text)
                logger.info(f"Text split into {len(chunks)} chunks (size: {chunk_size}, overlap: {chunk_overlap})")
                return chunks
                
        except Exception as e:
            logger.error(f"Error chunking text: {e}")
            raise
    
    def chunk_text_with_metadata(self, text: str, file_extension: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[Dict[str, Any]]:
        """텍스트를 청크로 분할하고 각 청크의 메타데이터(페이지, 라인 정보) 포함
        
        Args:
            text: 원본 텍스트
            file_extension: 파일 확장자
            chunk_size: 청크 크기
            chunk_overlap: 청크 간 겹침 크기
            
        Returns:
            청크별 메타데이터 리스트 [{"text": str, "page_number": int, "line_start": int, "line_end": int}, ...]
        """
        try:
            # 1. self.chunk_text로 청킹 수행
            chunks = self.chunk_text(text, chunk_size, chunk_overlap)
            
            # 2. 청크들을 합쳐서 재구성된 텍스트 생성 (구분자 제거된 상태)
            reconstructed_text = self._reconstruct_text_from_chunks(chunks, chunk_overlap)
            
            # 3. 재구성된 텍스트를 기준으로 메타데이터 테이블 구성
            line_table = self._build_line_offset_table(reconstructed_text, file_extension)
            
            # 4. 재구성된 텍스트에서 각 청크의 위치 찾기
            chunks_with_metadata = []
            current_pos = 0
            
            for idx, chunk in enumerate(chunks):
                # 재구성된 텍스트에서의 청크 위치
                chunk_start = current_pos
                chunk_end = current_pos + len(chunk) - 1
                
                # 라인 정보 추출 (재구성된 텍스트 기준)
                start_line_idx = self._find_line_index_by_pos(chunk_start, line_table)
                end_line_idx = self._find_line_index_by_pos(chunk_end, line_table)
                
                line_start = line_table[start_line_idx]["line_num"]
                line_end = line_table[end_line_idx]["line_num"]
                page_number = line_table[start_line_idx].get("page", 1)
                
                chunks_with_metadata.append({
                    "text": chunk,
                    "page_number": page_number,
                    "line_start": line_start,
                    "line_end": line_end,
                    "global_start": chunk_start,
                    "global_end": chunk_end,
                    "chunk_index": idx
                })
                
                # 다음 청크 위치로 이동
                current_pos += len(chunk)
                
                # 오버랩 처리 (마지막 청크가 아닌 경우)
                if idx < len(chunks) - 1:
                    next_chunk = chunks[idx + 1]
                    overlap_length = self._find_overlap_length(chunk, next_chunk, chunk_overlap)
                    current_pos -= overlap_length
            
            logger.info(f"Created {len(chunks_with_metadata)} chunks with metadata using reconstructed text")
            return chunks_with_metadata
            
        except Exception as e:
            logger.error(f"Failed to chunk text with metadata: {e}")
            # 실패 시 기본 청크만 반환
            chunks = self.chunk_text(text, chunk_size, chunk_overlap)
            return [{"text": chunk, "page_number": 1, "line_start": i+1, "line_end": i+1, "chunk_index": i} 
                for i, chunk in enumerate(chunks)]

    def _reconstruct_text_from_chunks(self, chunks: List[str], chunk_overlap: int) -> str:
        """청크들을 합쳐서 텍스트 재구성 (오버랩 제거)"""
        if not chunks:
            return ""
        
        if len(chunks) == 1:
            return chunks[0]
        
        reconstructed = chunks[0]
        
        for i in range(1, len(chunks)):
            current_chunk = chunks[i]
            prev_chunk = chunks[i-1]
            
            # 이전 청크와 현재 청크 간의 오버랩 길이 찾기
            overlap_length = self._find_overlap_length(prev_chunk, current_chunk, chunk_overlap)
            
            # 오버랩 부분을 제거하고 추가
            if overlap_length > 0:
                reconstructed += current_chunk[overlap_length:]
            else:
                reconstructed += current_chunk
        
        return reconstructed

    def _find_overlap_length(self, chunk1: str, chunk2: str, max_overlap: int) -> int:
        """두 청크 간의 실제 오버랩 길이 찾기"""
        max_check = min(len(chunk1), len(chunk2), max_overlap)
        
        for overlap_len in range(max_check, 0, -1):
            chunk1_suffix = chunk1[-overlap_len:]
            chunk2_prefix = chunk2[:overlap_len]
            
            if chunk1_suffix == chunk2_prefix:
                return overlap_len
        
        return 0
    
    def chunk_code_text(self, text: str, file_type: str, chunk_size: int = 1500, chunk_overlap: int = 300) -> List[str]:
        """코드 텍스트를 청크로 분할 (언어별 구문 구조를 고려한 분할)
            
        Args:
            text: 분할할 코드 텍스트
            file_type: 파일 형식
            chunk_size: 청크 크기 (코드는 좀 더 큰 청크 사용)
            chunk_overlap: 청크 간 중복 크기
            
        Returns:
            분할된 텍스트 청크 리스트
        """
        try:
            if not text or not text.strip():
                logger.warning("Empty code text provided for chunking")
                return [""]
            
            lang = LANGCHAIN_CODE_LANGUAGE_MAP.get(file_type.lower())

            if lang:
                text_splitter = RecursiveCharacterTextSplitter.from_language(
                    language=lang,
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap
                )
                logger.info(f"Using language-specific splitter for {file_type} ({lang})")
            else:
                # fallback: 기존 방식 사용
                logger.info(f"No language-specific splitter for {file_type}, using fallback.")
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap,
                    length_function=len,
                    separators=["\n\n", "\n", " ", ""]
                )

            chunks = text_splitter.split_text(text)
            logger.info(f"Code text split into {len(chunks)} chunks (size: {chunk_size}, overlap: {chunk_overlap})")
            return chunks
        except Exception as e:
            logger.error(f"Error chunking code text: {e}")
            raise
    
    def validate_file_format(self, file_path: str) -> tuple[bool, str]:
        """파일 형식 유효성 검사
        
        Args:
            file_path: 파일 경로
            
        Returns:
            (유효성, 파일형식) 튜플
        """
        try:
            file_extension = Path(file_path).suffix[1:].lower()
            is_valid = file_extension in self.supported_types
            return is_valid, file_extension
        except Exception:
            return False, ""
    
    def estimate_chunks_count(self, text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> int:
        """텍스트에서 생성될 청크 수 추정
        
        Args:
            text: 대상 텍스트
            chunk_size: 청크 크기
            chunk_overlap: 청크 간 중복 크기
            
        Returns:
            예상 청크 수
        """
        if not text:
            return 0
        
        text_length = len(text)
        if text_length <= chunk_size:
            return 1
        
        # 간단한 추정 공식
        effective_chunk_size = chunk_size - chunk_overlap
        estimated_chunks = (text_length - chunk_overlap) // effective_chunk_size + 1
        return max(1, estimated_chunks)
    
    def get_file_info(self, file_path: str) -> Dict[str, str]:
        """파일 정보 반환
        
        Args:
            file_path: 파일 경로
            
        Returns:
            파일 정보 딕셔너리 (extension, category, supported)
        """
        try:
            file_extension = Path(file_path).suffix[1:].lower()
            category = self.get_file_category(file_extension)
            is_supported = file_extension in self.supported_types
            
            return {
                'extension': file_extension,
                'category': category,
                'supported': str(is_supported)
            }
        except Exception:
            return {
                'extension': 'unknown',
                'category': 'unknown', 
                'supported': 'false'
            }

    def get_current_config_status(self) -> Dict[str, Any]:
        """현재 설정 상태 반환 (디버깅용)"""
        try:
            current_config = self._get_current_image_text_config()
            return {
                "provider": current_config.get('provider', 'unknown'),
                "ocr_enabled": self._is_image_text_enabled(current_config),
                "base_url": current_config.get('base_url', 'unknown'),
                "model": current_config.get('model', 'unknown'),
                "temperature": current_config.get('temperature', 'unknown'),
                "batch_size": current_config.get('batch_size', 1),
                "langchain_available": LANGCHAIN_OPENAI_AVAILABLE,
                "pdf2image_available": PDF2IMAGE_AVAILABLE,
                "docx2pdf_available": DOCX2PDF_AVAILABLE,
                "python_pptx_available": PYTHON_PPTX_AVAILABLE,
                "pil_available": PIL_AVAILABLE
            }
        except Exception as e:
            return {"error": str(e)}

    def test(self):
        """설정 테스트 메서드 (디버깅용)"""
        try:
            current_config = self._get_current_image_text_config()
            provider = current_config.get('provider', 'no_model')
            
            logger.info(f"🔍 Test - Current provider: {provider}")
            logger.info(f"🔍 Test - Current config: {current_config}")
            logger.info(f"🔍 Test - OCR enabled: {self._is_image_text_enabled(current_config)}")
            
        except Exception as e:
            logger.error(f"Error in test method: {e}")
            raise